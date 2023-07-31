#! python3
import datetime
import os
import sys
import time

# Work with DataFrames
import pandas as pd

# Work with Google Sheets
import pygsheets

# Work with MS Word files
from docx import Document
from docx.enum.section import WD_ORIENT
from docx.shared import Inches, Pt

# Work with MS PowerPoint files
from pptx import Presentation
from pptx.dml.color import RGBColor

# import PyQt5
from PyQt5 import QtCore, QtGui, QtWidgets

# Work with web browser automation
from selenium import webdriver
from selenium.common.exceptions import WebDriverException
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.common.by import By
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.support.ui import WebDriverWait

if hasattr(QtCore.Qt, "AA_EnableHighDpiScaling"):
    QtWidgets.QApplication.setAttribute(QtCore.Qt.AA_EnableHighDpiScaling, True)

if hasattr(QtCore.Qt, "AA+_UseHighDpiPixmaps"):
    QtWidgets.QApplication.setAttribute(QtCore.Qt.AA_UseHighDpiPixmaps, True)


# Main Window for GUI
class Ui_mainWindow(object):
    # Global variables and flags
    current_folder = os.path.dirname(os.path.realpath(__file__))
    GBCTemplate = os.path.join(current_folder, "Template-GBC.docx")
    GBCPptTemplate = os.path.join(current_folder, "Template-GBC.pptx")
    SFCTemplate = os.path.join(current_folder, "Template-SFC.docx")
    SFCPptTemplate = os.path.join(current_folder, "Template-SFC.pptx")
    genReport = False
    startDate = "2018-01-01"
    endDate = "2018-01-01"
    location = "Golden Bear Center"
    createSigns = False
    useExistingReport = False
    saveReportToPath = ""
    existingReportPath = ""
    classroomSignsOutput = False
    dailyScheduleOutput = False
    powerpointOutput = False
    saveSignsDirectory = ""
    uploadGBCSchedule = False
    GBCScheduleURL = ""
    uploadSFCSchedule = False
    SFCScheduleURL = ""
    center = {
        "Golden Bear Center": {
            "campus": "Berkeley - CA0001",
            "building": (
                "UC Berkeley Extension Golden Bear Center, 1995 University Ave. - GBC"
            ),
        },
        "San Francisco Center": {
            "campus": "San Francisco - CA0003",
            "building": "San Francisco Campus, 160 Spear St. - SFCAMPUS",
        },
    }

    centerReverse = {
        "GBC - UC Berkeley Extension Golden Bear Center, 1995 University Ave.": {
            "name": "GBC",
            "template": GBCTemplate,
            "pptTemplate": GBCPptTemplate,
        },
        "SFCAMPUS - San Francisco Campus, 160 Spear St.": {
            "name": "SFC",
            "template": SFCTemplate,
            "pptTemplate": SFCPptTemplate,
        },
    }

    def __init__(self) -> None:
        super().__init__()
        # Initialize settings from config.ini file, otherwise set default
        self.settings = QtCore.QSettings("config.ini", QtCore.QSettings.IniFormat)
        self.genReport = self.settings.value("genReport", True, type=bool)
        self.saveReportToPath = self.settings.value(
            "saveReportToPath", self.current_folder, type=str
        )
        self.createSigns = self.settings.value("createSigns", False, type=bool)
        self.useExistingReport = self.settings.value(
            "useExistingReport", False, type=bool
        )
        self.classroomSignsOutput = self.settings.value(
            "classroomSignsOutput", False, type=bool
        )
        self.dailyScheduleOutput = self.settings.value(
            "dailyScheduleOutput", False, type=bool
        )
        self.powerpointOutput = self.settings.value(
            "powerpointOutput", False, type=bool
        )
        self.saveSignsDirectory = self.settings.value(
            "saveSignsDirectory", self.saveReportToPath, type=str
        )
        self.uploadGBCSchedule = self.settings.value(
            "uploadGBCSchedule", True, type=bool
        )
        self.uploadSFCSchedule = self.settings.value(
            "uploadSFCSchedule", True, type=bool
        )
        self.GBCScheduleURL = self.settings.value("GBCScheduleURL", "", type=str)
        self.SFCScheduleURL = self.settings.value("SFCScheduleURL", "", type=str)

    def setupUi(self, mainWindow: QtWidgets.QWidget) -> None:
        # global startDate, endDate
        mainWindow.setObjectName("mainWindow")
        mainWindow.setWindowModality(QtCore.Qt.NonModal)
        mainWindow.setEnabled(True)
        mainWindow.resize(547, 425)
        sizePolicy = QtWidgets.QSizePolicy(
            QtWidgets.QSizePolicy.Maximum, QtWidgets.QSizePolicy.Maximum
        )
        sizePolicy.setHorizontalStretch(0)
        sizePolicy.setVerticalStretch(0)
        sizePolicy.setHeightForWidth(mainWindow.sizePolicy().hasHeightForWidth())
        mainWindow.setSizePolicy(sizePolicy)
        mainWindow.setBaseSize(QtCore.QSize(430, 400))
        self.mainWindowLayout = QtWidgets.QVBoxLayout(mainWindow)
        self.mainWindowLayout.setObjectName("mainWindowLayout")
        self.genReportBox = QtWidgets.QGroupBox(mainWindow)
        sizePolicy = QtWidgets.QSizePolicy(
            QtWidgets.QSizePolicy.Expanding, QtWidgets.QSizePolicy.MinimumExpanding
        )
        sizePolicy.setHorizontalStretch(0)
        sizePolicy.setVerticalStretch(0)
        sizePolicy.setHeightForWidth(self.genReportBox.sizePolicy().hasHeightForWidth())
        self.genReportBox.setSizePolicy(sizePolicy)
        font = QtGui.QFont()
        font.setFamily("Tahoma")
        font.setPointSize(14)
        font.setBold(True)
        font.setItalic(False)
        font.setUnderline(True)
        font.setWeight(75)
        self.genReportBox.setFont(font)
        self.genReportBox.setCheckable(True)
        self.genReportBox.setChecked(self.genReport)
        self.genReportBox.setObjectName("genReportBox")
        self.genReportLayout = QtWidgets.QVBoxLayout(self.genReportBox)
        self.genReportLayout.setSizeConstraint(QtWidgets.QLayout.SetDefaultConstraint)
        self.genReportLayout.setObjectName("genReportLayout")
        self.dateLayout = QtWidgets.QHBoxLayout()
        self.dateLayout.setSizeConstraint(QtWidgets.QLayout.SetDefaultConstraint)
        self.dateLayout.setObjectName("dateLayout")
        self.startDateLabel = QtWidgets.QLabel(self.genReportBox)
        sizePolicy = QtWidgets.QSizePolicy(
            QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Minimum
        )
        sizePolicy.setHorizontalStretch(0)
        sizePolicy.setVerticalStretch(0)
        sizePolicy.setHeightForWidth(
            self.startDateLabel.sizePolicy().hasHeightForWidth()
        )
        self.startDateLabel.setSizePolicy(sizePolicy)
        font = QtGui.QFont()
        font.setFamily("Tahoma")
        font.setPointSize(14)
        font.setBold(True)
        font.setItalic(False)
        font.setUnderline(False)
        font.setWeight(75)
        self.startDateLabel.setFont(font)
        self.startDateLabel.setObjectName("startDateLabel")
        self.dateLayout.addWidget(self.startDateLabel)
        self.selectStartDate = QtWidgets.QDateEdit(self.genReportBox)
        sizePolicy = QtWidgets.QSizePolicy(
            QtWidgets.QSizePolicy.MinimumExpanding, QtWidgets.QSizePolicy.Minimum
        )
        sizePolicy.setHorizontalStretch(0)
        sizePolicy.setVerticalStretch(0)
        sizePolicy.setHeightForWidth(
            self.selectStartDate.sizePolicy().hasHeightForWidth()
        )
        self.selectStartDate.setSizePolicy(sizePolicy)
        font = QtGui.QFont()
        font.setFamily("Tahoma")
        font.setPointSize(14)
        font.setBold(True)
        font.setItalic(False)
        font.setUnderline(False)
        font.setWeight(75)
        self.selectStartDate.setFont(font)
        self.selectStartDate.setFrame(True)
        self.selectStartDate.setReadOnly(False)
        self.selectStartDate.setProperty("showGroupSeparator", False)
        self.selectStartDate.setCalendarPopup(True)
        self.selectStartDate.setDate(QtCore.QDate.currentDate())
        self.startDate = str(QtCore.QDate.currentDate().toPyDate())
        self.selectStartDate.setObjectName("selectStartDate")
        self.dateLayout.addWidget(self.selectStartDate)
        self.endDateLabel = QtWidgets.QLabel(self.genReportBox)
        sizePolicy = QtWidgets.QSizePolicy(
            QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Minimum
        )
        sizePolicy.setHorizontalStretch(0)
        sizePolicy.setVerticalStretch(0)
        sizePolicy.setHeightForWidth(self.endDateLabel.sizePolicy().hasHeightForWidth())
        self.endDateLabel.setSizePolicy(sizePolicy)
        font = QtGui.QFont()
        font.setFamily("Tahoma")
        font.setPointSize(14)
        font.setBold(True)
        font.setItalic(False)
        font.setUnderline(False)
        font.setWeight(75)
        self.endDateLabel.setFont(font)
        self.endDateLabel.setObjectName("endDateLabel")
        self.dateLayout.addWidget(self.endDateLabel)
        self.selectEndDate = QtWidgets.QDateEdit(self.genReportBox)
        sizePolicy = QtWidgets.QSizePolicy(
            QtWidgets.QSizePolicy.MinimumExpanding, QtWidgets.QSizePolicy.Minimum
        )
        sizePolicy.setHorizontalStretch(0)
        sizePolicy.setVerticalStretch(0)
        sizePolicy.setHeightForWidth(
            self.selectEndDate.sizePolicy().hasHeightForWidth()
        )
        self.selectEndDate.setSizePolicy(sizePolicy)
        font = QtGui.QFont()
        font.setFamily("Tahoma")
        font.setPointSize(14)
        font.setBold(True)
        font.setItalic(False)
        font.setUnderline(False)
        font.setWeight(75)
        self.selectEndDate.setFont(font)
        self.selectEndDate.setFocusPolicy(QtCore.Qt.WheelFocus)
        self.selectEndDate.setReadOnly(False)
        self.selectEndDate.setCalendarPopup(True)
        self.selectEndDate.setDate(QtCore.QDate.currentDate())
        self.endDate = str(QtCore.QDate.currentDate().toPyDate())
        self.selectEndDate.setObjectName("selectEndDate")
        self.dateLayout.addWidget(self.selectEndDate)
        self.genReportLayout.addLayout(self.dateLayout)
        self.locationLayout = QtWidgets.QHBoxLayout()
        self.locationLayout.setObjectName("locationLayout")
        self.locationLabel = QtWidgets.QLabel(self.genReportBox)
        sizePolicy = QtWidgets.QSizePolicy(
            QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Minimum
        )
        sizePolicy.setHorizontalStretch(0)
        sizePolicy.setVerticalStretch(0)
        sizePolicy.setHeightForWidth(
            self.locationLabel.sizePolicy().hasHeightForWidth()
        )
        self.locationLabel.setSizePolicy(sizePolicy)
        font = QtGui.QFont()
        font.setFamily("Tahoma")
        font.setPointSize(14)
        font.setBold(True)
        font.setItalic(False)
        font.setUnderline(False)
        font.setWeight(75)
        self.locationLabel.setFont(font)
        self.locationLabel.setObjectName("locationLabel")
        self.locationLayout.addWidget(self.locationLabel)
        self.selectLocation = QtWidgets.QComboBox(self.genReportBox)
        sizePolicy = QtWidgets.QSizePolicy(
            QtWidgets.QSizePolicy.MinimumExpanding, QtWidgets.QSizePolicy.Minimum
        )
        sizePolicy.setHorizontalStretch(0)
        sizePolicy.setVerticalStretch(0)
        sizePolicy.setHeightForWidth(
            self.selectLocation.sizePolicy().hasHeightForWidth()
        )
        self.selectLocation.setSizePolicy(sizePolicy)
        font = QtGui.QFont()
        font.setFamily("Tahoma")
        font.setPointSize(14)
        font.setBold(True)
        font.setItalic(False)
        font.setUnderline(False)
        font.setWeight(75)
        self.selectLocation.setFont(font)
        self.selectLocation.setEditable(False)
        self.selectLocation.setObjectName("selectLocation")
        self.selectLocation.addItem("")
        self.selectLocation.addItem("")
        self.locationLayout.addWidget(self.selectLocation)
        self.genReportLayout.addLayout(self.locationLayout)
        self.saveReportPathLayout = QtWidgets.QHBoxLayout()
        self.saveReportPathLayout.setObjectName("saveReportPathLayout")
        self.saveReportPathLabel = QtWidgets.QLabel(self.genReportBox)
        sizePolicy = QtWidgets.QSizePolicy(
            QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Minimum
        )
        sizePolicy.setHorizontalStretch(0)
        sizePolicy.setVerticalStretch(0)
        sizePolicy.setHeightForWidth(
            self.saveReportPathLabel.sizePolicy().hasHeightForWidth()
        )
        self.saveReportPathLabel.setSizePolicy(sizePolicy)
        font = QtGui.QFont()
        font.setFamily("Tahoma")
        font.setPointSize(14)
        font.setBold(True)
        font.setItalic(False)
        font.setUnderline(False)
        font.setWeight(75)
        self.saveReportPathLabel.setFont(font)
        self.saveReportPathLabel.setObjectName("saveReportPathLabel")
        self.saveReportPathLayout.addWidget(self.saveReportPathLabel)
        self.selectSaveReportPath = QtWidgets.QLineEdit(self.genReportBox)
        sizePolicy = QtWidgets.QSizePolicy(
            QtWidgets.QSizePolicy.MinimumExpanding, QtWidgets.QSizePolicy.Minimum
        )
        sizePolicy.setHorizontalStretch(0)
        sizePolicy.setVerticalStretch(0)
        sizePolicy.setHeightForWidth(
            self.selectSaveReportPath.sizePolicy().hasHeightForWidth()
        )
        self.selectSaveReportPath.setSizePolicy(sizePolicy)
        font = QtGui.QFont()
        font.setFamily("Tahoma")
        font.setPointSize(14)
        font.setBold(True)
        font.setItalic(False)
        font.setUnderline(False)
        font.setWeight(75)
        self.selectSaveReportPath.setFont(font)
        self.selectSaveReportPath.setReadOnly(True)
        self.selectSaveReportPath.setObjectName("selectSaveReportPath")
        self.saveReportPathLayout.addWidget(self.selectSaveReportPath)
        self.browseSaveReportButton = QtWidgets.QToolButton(self.genReportBox)
        sizePolicy = QtWidgets.QSizePolicy(
            QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Minimum
        )
        sizePolicy.setHorizontalStretch(0)
        sizePolicy.setVerticalStretch(0)
        sizePolicy.setHeightForWidth(
            self.browseSaveReportButton.sizePolicy().hasHeightForWidth()
        )
        self.browseSaveReportButton.setSizePolicy(sizePolicy)
        font = QtGui.QFont()
        font.setFamily("Tahoma")
        font.setPointSize(14)
        font.setBold(True)
        font.setItalic(False)
        font.setUnderline(False)
        font.setWeight(75)
        self.browseSaveReportButton.setFont(font)
        self.browseSaveReportButton.setCheckable(False)
        self.browseSaveReportButton.setPopupMode(QtWidgets.QToolButton.InstantPopup)
        self.browseSaveReportButton.setToolButtonStyle(QtCore.Qt.ToolButtonTextOnly)
        self.browseSaveReportButton.setObjectName("browseSaveReportButton")
        self.saveReportPathLayout.addWidget(self.browseSaveReportButton)
        self.genReportLayout.addLayout(self.saveReportPathLayout)
        self.mainWindowLayout.addWidget(self.genReportBox)
        self.createSignsBox = QtWidgets.QGroupBox(mainWindow)
        self.createSignsBox.setEnabled(True)
        sizePolicy = QtWidgets.QSizePolicy(
            QtWidgets.QSizePolicy.Expanding, QtWidgets.QSizePolicy.MinimumExpanding
        )
        sizePolicy.setHorizontalStretch(0)
        sizePolicy.setVerticalStretch(0)
        sizePolicy.setHeightForWidth(
            self.createSignsBox.sizePolicy().hasHeightForWidth()
        )
        self.createSignsBox.setSizePolicy(sizePolicy)
        self.createSignsBox.setMinimumSize(QtCore.QSize(150, 194))
        font = QtGui.QFont()
        font.setFamily("Tahoma")
        font.setPointSize(14)
        font.setBold(True)
        font.setItalic(False)
        font.setUnderline(True)
        font.setWeight(75)
        self.createSignsBox.setFont(font)
        self.createSignsBox.setCheckable(True)
        self.createSignsBox.setChecked(self.createSigns)
        self.createSignsBox.setObjectName("createSignsBox")
        self.createSignsBoxLayout = QtWidgets.QVBoxLayout(self.createSignsBox)
        self.createSignsBoxLayout.setObjectName("createSignsBoxLayout")
        self.useExistingReportBox = QtWidgets.QGroupBox(self.createSignsBox)
        sizePolicy = QtWidgets.QSizePolicy(
            QtWidgets.QSizePolicy.MinimumExpanding, QtWidgets.QSizePolicy.Minimum
        )
        sizePolicy.setHorizontalStretch(0)
        sizePolicy.setVerticalStretch(0)
        sizePolicy.setHeightForWidth(
            self.useExistingReportBox.sizePolicy().hasHeightForWidth()
        )
        self.useExistingReportBox.setSizePolicy(sizePolicy)
        self.useExistingReportBox.setMinimumSize(QtCore.QSize(0, 76))
        font = QtGui.QFont()
        font.setFamily("Tahoma")
        font.setPointSize(14)
        font.setBold(True)
        font.setItalic(False)
        font.setUnderline(False)
        font.setWeight(75)
        self.useExistingReportBox.setFont(font)
        self.useExistingReportBox.setCheckable(True)
        self.useExistingReportBox.setChecked(self.useExistingReport)
        self.useExistingReportBox.setObjectName("useExistingReportBox")
        self.useExistingReportLayout = QtWidgets.QVBoxLayout(self.useExistingReportBox)
        self.useExistingReportLayout.setObjectName("useExistingReportLayout")
        self.browseExistingReportLayout = QtWidgets.QHBoxLayout()
        self.browseExistingReportLayout.setObjectName("browseExistingReportLayout")
        self.selectExistingReportPath = QtWidgets.QLineEdit(self.useExistingReportBox)
        sizePolicy = QtWidgets.QSizePolicy(
            QtWidgets.QSizePolicy.MinimumExpanding, QtWidgets.QSizePolicy.Minimum
        )
        sizePolicy.setHorizontalStretch(0)
        sizePolicy.setVerticalStretch(0)
        sizePolicy.setHeightForWidth(
            self.selectExistingReportPath.sizePolicy().hasHeightForWidth()
        )
        self.selectExistingReportPath.setSizePolicy(sizePolicy)
        self.selectExistingReportPath.setMinimumSize(QtCore.QSize(0, 30))
        font = QtGui.QFont()
        font.setFamily("Tahoma")
        font.setPointSize(14)
        font.setBold(True)
        font.setItalic(False)
        font.setWeight(75)
        self.selectExistingReportPath.setFont(font)
        self.selectExistingReportPath.setReadOnly(True)
        self.selectExistingReportPath.setObjectName("selectExistingReportPath")
        self.browseExistingReportLayout.addWidget(self.selectExistingReportPath)
        self.browseExistingReportButton = QtWidgets.QToolButton(
            self.useExistingReportBox
        )
        sizePolicy = QtWidgets.QSizePolicy(
            QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Minimum
        )
        sizePolicy.setHorizontalStretch(0)
        sizePolicy.setVerticalStretch(0)
        sizePolicy.setHeightForWidth(
            self.browseExistingReportButton.sizePolicy().hasHeightForWidth()
        )
        self.browseExistingReportButton.setSizePolicy(sizePolicy)
        font = QtGui.QFont()
        font.setFamily("Tahoma")
        font.setPointSize(14)
        font.setBold(True)
        font.setItalic(False)
        font.setWeight(75)
        self.browseExistingReportButton.setFont(font)
        self.browseExistingReportButton.setPopupMode(QtWidgets.QToolButton.InstantPopup)
        self.browseExistingReportButton.setToolButtonStyle(QtCore.Qt.ToolButtonTextOnly)
        self.browseExistingReportButton.setAutoRaise(False)
        self.browseExistingReportButton.setArrowType(QtCore.Qt.NoArrow)
        self.browseExistingReportButton.setObjectName("browseExistingReportButton")
        self.browseExistingReportLayout.addWidget(self.browseExistingReportButton)
        self.useExistingReportLayout.addLayout(self.browseExistingReportLayout)
        self.createSignsBoxLayout.addWidget(self.useExistingReportBox)
        self.saveSignsPathLayout = QtWidgets.QHBoxLayout()
        self.saveSignsPathLayout.setObjectName("saveSignsPathLayout")
        self.saveSignsPathLabel = QtWidgets.QLabel(self.createSignsBox)
        sizePolicy = QtWidgets.QSizePolicy(
            QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Minimum
        )
        sizePolicy.setHorizontalStretch(0)
        sizePolicy.setVerticalStretch(0)
        sizePolicy.setHeightForWidth(
            self.saveSignsPathLabel.sizePolicy().hasHeightForWidth()
        )
        self.saveSignsPathLabel.setSizePolicy(sizePolicy)
        self.saveSignsPathLabel.setMinimumSize(QtCore.QSize(101, 0))
        font = QtGui.QFont()
        font.setFamily("Tahoma")
        font.setPointSize(14)
        font.setBold(True)
        font.setItalic(False)
        font.setUnderline(False)
        font.setWeight(75)
        self.saveSignsPathLabel.setFont(font)
        self.saveSignsPathLabel.setObjectName("saveSignsPathLabel")
        self.saveSignsPathLayout.addWidget(self.saveSignsPathLabel)
        self.selectSaveSignsPath = QtWidgets.QLineEdit(self.createSignsBox)
        sizePolicy = QtWidgets.QSizePolicy(
            QtWidgets.QSizePolicy.MinimumExpanding, QtWidgets.QSizePolicy.Minimum
        )
        sizePolicy.setHorizontalStretch(0)
        sizePolicy.setVerticalStretch(0)
        sizePolicy.setHeightForWidth(
            self.selectSaveSignsPath.sizePolicy().hasHeightForWidth()
        )
        self.selectSaveSignsPath.setSizePolicy(sizePolicy)
        self.selectSaveSignsPath.setMinimumSize(QtCore.QSize(304, 0))
        font = QtGui.QFont()
        font.setFamily("Tahoma")
        font.setPointSize(14)
        font.setBold(True)
        font.setItalic(False)
        font.setUnderline(False)
        font.setWeight(75)
        self.selectSaveSignsPath.setFont(font)
        self.selectSaveSignsPath.setReadOnly(True)
        self.selectSaveSignsPath.setObjectName("selectSaveSignsPath")
        self.saveSignsPathLayout.addWidget(self.selectSaveSignsPath)
        self.browseSaveSignsButton = QtWidgets.QToolButton(self.createSignsBox)
        sizePolicy = QtWidgets.QSizePolicy(
            QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Minimum
        )
        sizePolicy.setHorizontalStretch(0)
        sizePolicy.setVerticalStretch(0)
        sizePolicy.setHeightForWidth(
            self.browseSaveSignsButton.sizePolicy().hasHeightForWidth()
        )
        self.browseSaveSignsButton.setSizePolicy(sizePolicy)
        self.browseSaveSignsButton.setMinimumSize(QtCore.QSize(90, 0))
        font = QtGui.QFont()
        font.setFamily("Tahoma")
        font.setPointSize(14)
        font.setBold(True)
        font.setItalic(False)
        font.setUnderline(False)
        font.setWeight(75)
        self.browseSaveSignsButton.setFont(font)
        self.browseSaveSignsButton.setCheckable(False)
        self.browseSaveSignsButton.setPopupMode(QtWidgets.QToolButton.InstantPopup)
        self.browseSaveSignsButton.setToolButtonStyle(QtCore.Qt.ToolButtonTextOnly)
        self.browseSaveSignsButton.setObjectName("browseSaveSignsButton")
        self.saveSignsPathLayout.addWidget(self.browseSaveSignsButton)
        self.createSignsBoxLayout.addLayout(self.saveSignsPathLayout)
        self.outputOptionsLayout = QtWidgets.QHBoxLayout()
        self.outputOptionsLayout.setSpacing(12)
        self.outputOptionsLayout.setObjectName("outputOptionsLayout")
        self.classroomSignsCheckbox = QtWidgets.QCheckBox(self.createSignsBox)
        sizePolicy = QtWidgets.QSizePolicy(
            QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Minimum
        )
        sizePolicy.setHorizontalStretch(0)
        sizePolicy.setVerticalStretch(0)
        sizePolicy.setHeightForWidth(
            self.classroomSignsCheckbox.sizePolicy().hasHeightForWidth()
        )
        self.classroomSignsCheckbox.setSizePolicy(sizePolicy)
        self.classroomSignsCheckbox.setMinimumSize(QtCore.QSize(0, 0))
        font = QtGui.QFont()
        font.setFamily("Tahoma")
        font.setPointSize(14)
        font.setBold(True)
        font.setItalic(False)
        font.setUnderline(False)
        font.setWeight(75)
        self.classroomSignsCheckbox.setFont(font)
        self.classroomSignsCheckbox.setChecked(self.classroomSignsOutput)
        self.classroomSignsCheckbox.setObjectName("classroomSignsCheckbox")
        self.outputOptionsLayout.addWidget(self.classroomSignsCheckbox)
        self.dailyScheduleCheckbox = QtWidgets.QCheckBox(self.createSignsBox)
        sizePolicy = QtWidgets.QSizePolicy(
            QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Minimum
        )
        sizePolicy.setHorizontalStretch(0)
        sizePolicy.setVerticalStretch(0)
        sizePolicy.setHeightForWidth(
            self.dailyScheduleCheckbox.sizePolicy().hasHeightForWidth()
        )
        self.dailyScheduleCheckbox.setSizePolicy(sizePolicy)
        self.dailyScheduleCheckbox.setMinimumSize(QtCore.QSize(0, 0))
        font = QtGui.QFont()
        font.setFamily("Tahoma")
        font.setPointSize(14)
        font.setBold(True)
        font.setItalic(False)
        font.setUnderline(False)
        font.setWeight(75)
        self.dailyScheduleCheckbox.setFont(font)
        self.dailyScheduleCheckbox.setChecked(self.dailyScheduleOutput)
        self.dailyScheduleCheckbox.setObjectName("dailyScheduleCheckbox")
        self.outputOptionsLayout.addWidget(self.dailyScheduleCheckbox)
        self.powerpointCheckbox = QtWidgets.QCheckBox(self.createSignsBox)
        sizePolicy = QtWidgets.QSizePolicy(
            QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Minimum
        )
        sizePolicy.setHorizontalStretch(0)
        sizePolicy.setVerticalStretch(0)
        sizePolicy.setHeightForWidth(
            self.powerpointCheckbox.sizePolicy().hasHeightForWidth()
        )
        self.powerpointCheckbox.setSizePolicy(sizePolicy)
        self.powerpointCheckbox.setMinimumSize(QtCore.QSize(0, 0))
        font = QtGui.QFont()
        font.setFamily("Tahoma")
        font.setPointSize(14)
        font.setBold(True)
        font.setItalic(False)
        font.setUnderline(False)
        font.setWeight(75)
        self.powerpointCheckbox.setFont(font)
        self.powerpointCheckbox.setChecked(self.powerpointOutput)
        self.powerpointCheckbox.setObjectName("powerpointCheckbox")
        self.outputOptionsLayout.addWidget(self.powerpointCheckbox)
        self.createSignsBoxLayout.addLayout(self.outputOptionsLayout)
        self.mainWindowLayout.addWidget(self.createSignsBox)
        self.startExitLayout = QtWidgets.QHBoxLayout()
        self.startExitLayout.setSizeConstraint(QtWidgets.QLayout.SetDefaultConstraint)
        self.startExitLayout.setObjectName("startExitLayout")
        spacerItem = QtWidgets.QSpacerItem(
            40, 20, QtWidgets.QSizePolicy.Expanding, QtWidgets.QSizePolicy.Minimum
        )
        self.startExitLayout.addItem(spacerItem)
        self.StartButton = QtWidgets.QPushButton(mainWindow)
        sizePolicy = QtWidgets.QSizePolicy(
            QtWidgets.QSizePolicy.MinimumExpanding, QtWidgets.QSizePolicy.Minimum
        )
        sizePolicy.setHorizontalStretch(0)
        sizePolicy.setVerticalStretch(0)
        sizePolicy.setHeightForWidth(self.StartButton.sizePolicy().hasHeightForWidth())
        self.StartButton.setSizePolicy(sizePolicy)
        self.StartButton.setMinimumSize(QtCore.QSize(175, 50))
        font = QtGui.QFont()
        font.setFamily("Tahoma")
        font.setPointSize(14)
        font.setBold(True)
        font.setItalic(False)
        font.setWeight(75)
        self.StartButton.setFont(font)
        self.StartButton.setObjectName("StartButton")
        self.startExitLayout.addWidget(self.StartButton)
        spacerItem1 = QtWidgets.QSpacerItem(
            40,
            20,
            QtWidgets.QSizePolicy.MinimumExpanding,
            QtWidgets.QSizePolicy.Minimum,
        )
        self.startExitLayout.addItem(spacerItem1)
        self.exitButton = QtWidgets.QPushButton(mainWindow)
        sizePolicy = QtWidgets.QSizePolicy(
            QtWidgets.QSizePolicy.MinimumExpanding, QtWidgets.QSizePolicy.Minimum
        )
        sizePolicy.setHorizontalStretch(0)
        sizePolicy.setVerticalStretch(0)
        sizePolicy.setHeightForWidth(self.exitButton.sizePolicy().hasHeightForWidth())
        self.exitButton.setSizePolicy(sizePolicy)
        self.exitButton.setMinimumSize(QtCore.QSize(175, 50))
        font = QtGui.QFont()
        font.setFamily("Tahoma")
        font.setPointSize(14)
        font.setBold(True)
        font.setItalic(False)
        font.setWeight(75)
        self.exitButton.setFont(font)
        self.exitButton.setObjectName("exitButton")
        self.startExitLayout.addWidget(self.exitButton)
        spacerItem2 = QtWidgets.QSpacerItem(
            40, 20, QtWidgets.QSizePolicy.Expanding, QtWidgets.QSizePolicy.Minimum
        )
        self.startExitLayout.addItem(spacerItem2)
        self.mainWindowLayout.addLayout(self.startExitLayout)

        self.retranslateUi(mainWindow)
        QtCore.QMetaObject.connectSlotsByName(mainWindow)

    def retranslateUi(self, mainWindow: QtWidgets.QWidget) -> None:
        _translate = QtCore.QCoreApplication.translate
        mainWindow.setWindowTitle(_translate("mainWindow", "Auto Schedule"))
        self.genReportBox.setTitle(_translate("mainWindow", "Generate Destiny Report"))
        self.startDateLabel.setText(_translate("mainWindow", "Start Date:"))
        self.endDateLabel.setText(_translate("mainWindow", "End Date:"))
        self.locationLabel.setText(_translate("mainWindow", "Location:"))
        self.selectLocation.setItemText(
            0, _translate("mainWindow", "Golden Bear Center")
        )
        self.selectLocation.setItemText(
            1, _translate("mainWindow", "San Francisco Center")
        )
        self.saveReportPathLabel.setText(_translate("mainWindow", "Save Path:"))
        # self.selectSaveReportPath.setText(_translate("mainWindow", ""))
        self.selectSaveReportPath.setText(
            _translate("mainWindow", self.saveReportToPath)
        )
        self.browseSaveReportButton.setText(_translate("mainWindow", "Browse"))
        self.createSignsBox.setTitle(_translate("mainWindow", "Create Signs"))
        self.useExistingReportBox.setTitle(
            _translate("mainWindow", "Use existing Destiny Report")
        )
        self.selectExistingReportPath.setText(_translate("mainWindow", ""))
        self.browseExistingReportButton.setText(_translate("mainWindow", "Browse"))
        self.saveSignsPathLabel.setText(_translate("mainWindow", "Save Path:"))
        self.selectSaveSignsPath.setText(
            _translate("mainWindow", self.saveSignsDirectory)
        )
        self.browseSaveSignsButton.setText(_translate("mainWindow", "Browse"))
        self.classroomSignsCheckbox.setText(_translate("mainWindow", "Classroom Sign"))
        self.dailyScheduleCheckbox.setText(_translate("mainWindow", "Daily Schedule"))
        self.powerpointCheckbox.setText(_translate("mainWindow", "PowerPoint"))
        self.StartButton.setText(_translate("mainWindow", "Start"))
        self.exitButton.setText(_translate("mainWindow", "Exit"))

        self.genReportBox.toggled.connect(self.genReportState)
        self.selectStartDate.dateChanged.connect(self.startDateChanged)
        self.selectEndDate.dateChanged.connect(self.endDateChanged)
        self.selectLocation.currentIndexChanged.connect(self.locationChanged)
        self.browseSaveReportButton.clicked.connect(self.saveReportDirectory)
        self.createSignsBox.toggled.connect(self.createSignsState)
        self.useExistingReportBox.toggled.connect(self.useExistingReportState)
        self.browseExistingReportButton.clicked.connect(self.browseExistingReport)
        self.classroomSignsCheckbox.toggled.connect(self.classroomSignsState)
        self.dailyScheduleCheckbox.toggled.connect(self.dailyScheduleState)
        self.powerpointCheckbox.toggled.connect(self.powerpointState)
        self.browseSaveSignsButton.clicked.connect(self.saveSignsPath)
        self.exitButton.clicked.connect(self.exitApp)
        self.StartButton.clicked.connect(self.startApp)

    def genReportState(self) -> None:
        if self.genReportBox.isChecked():
            self.genReport = True
            self.useExistingReport = False
            self.useExistingReportBox.setChecked(False)
        else:
            self.genReport = False
            self.useExistingReport = True
            self.useExistingReportBox.setEnabled(True)
            self.useExistingReportBox.setChecked(True)

    def startDateChanged(self) -> None:
        self.startDate = str(self.selectStartDate.date().toPyDate())
        self.endDate = self.startDate
        self.selectEndDate.setDate(self.selectStartDate.date())

    def endDateChanged(self) -> None:
        self.endDate = str(self.selectEndDate.date().toPyDate())

    def locationChanged(self) -> None:
        self.location = self.selectLocation.currentText()

    def saveReportDirectory(self) -> None:
        path = os.path.normpath(
            QtWidgets.QFileDialog.getExistingDirectory(
                None, "Save Destiny Report to", self.saveReportToPath
            )
        )
        if path and path != ".":
            self.saveReportToPath = path
            self.selectSaveReportPath.setText(self.saveReportToPath)

    def createSignsState(self) -> None:
        if self.createSignsBox.isChecked():
            self.createSigns = True
            if not self.genReportBox.isChecked():
                self.useExistingReport = True
                self.useExistingReportBox.setChecked(True)
                self.useExistingReportBox.setEnabled(True)
        else:
            self.createSigns = False
            self.useExistingReport = False
            self.useExistingReportBox.setChecked(False)
            self.useExistingReportBox.setEnabled(False)

    def useExistingReportState(self) -> None:
        if self.useExistingReportBox.isChecked():
            self.useExistingReport = True
            self.genReport = False
            self.genReportBox.setChecked(False)
        else:
            self.useExistingReport = False
            self.genReport = True
            self.genReportBox.setChecked(True)

    def browseExistingReport(self) -> None:
        path, _ = QtWidgets.QFileDialog.getOpenFileName(
            None,
            "Select SectionScheduleDailySummary.xls",
            self.saveReportToPath,
            "Excel Files (*.xls)",
        )
        if path and path != ".":
            self.existingReportPath = os.path.normpath(path)
            self.selectExistingReportPath.setText(self.existingReportPath)

    def classroomSignsState(self) -> None:
        if self.classroomSignsCheckbox.isChecked():
            self.classroomSignsOutput = True
        else:
            self.classroomSignsOutput = False

    def dailyScheduleState(self) -> None:
        if self.dailyScheduleCheckbox.isChecked():
            self.dailyScheduleOutput = True
        else:
            self.dailyScheduleOutput = False

    def powerpointState(self) -> None:
        if self.powerpointCheckbox.isChecked():
            self.powerpointOutput = True
        else:
            self.powerpointOutput = False

    def saveSignsPath(self) -> None:
        path = os.path.normpath(
            QtWidgets.QFileDialog.getExistingDirectory(
                None, "Save Signs to", self.saveSignsDirectory
            )
        )
        if path and path != ".":
            self.saveSignsDirectory = path
            self.selectSaveSignsPath.setText(self.saveSignsDirectory)

    def exitApp(self) -> None:
        reply = QtWidgets.QMessageBox.question(
            None,
            "Exit",
            "Are you sure you want to exit?",
            QtWidgets.QMessageBox.Yes | QtWidgets.QMessageBox.No,
            QtWidgets.QMessageBox.No,
        )
        if reply == QtWidgets.QMessageBox.Yes:
            # Save settings state to config.ini file
            self.settings.setValue("genReport", self.genReport)
            self.settings.setValue("saveReportToPath", self.saveReportToPath)
            self.settings.setValue("createSigns", self.createSigns)
            self.settings.setValue("useExistingReport", self.useExistingReport)
            self.settings.setValue("classroomSignsOutput", self.classroomSignsOutput)
            self.settings.setValue("dailyScheduleOutput", self.dailyScheduleOutput)
            self.settings.setValue("powerpointOutput", self.powerpointOutput)
            self.settings.setValue("saveSignsDirectory", self.saveSignsDirectory)
            self.settings.setValue("uploadGBCSchedule", self.uploadGBCSchedule)
            self.settings.setValue("GBCScheduleURL", self.GBCScheduleURL)
            self.settings.setValue("uploadSFCSchedule", self.uploadSFCSchedule)
            self.settings.setValue("SFCScheduleURL", self.SFCScheduleURL)
            sys.exit()
        else:
            pass

    def startApp(self) -> None:
        # Checks for invalid user input and all required input is selected.
        result = 1
        if not self.genReport and not self.createSigns:
            QtWidgets.QMessageBox.warning(
                None, "No settings", "Please fill in the settings first."
            )
            return
        if self.genReport:
            if self.endDate < self.startDate:
                QtWidgets.QMessageBox.warning(
                    None, "Invalid date range", "Please select a valid date range."
                )
                return
            elif self.saveReportToPath == "":
                QtWidgets.QMessageBox.warning(
                    None,
                    "Save location error",
                    "Please select where you want to save the report to.",
                )
                return
            else:
                if self.createSigns:
                    if self.saveSignsDirectory == "" or not os.path.isdir(
                        self.saveSignsDirectory
                    ):
                        QtWidgets.QMessageBox.warning(
                            None,
                            "Save location error",
                            "Please select where you want to save the signs to.",
                        )
                        return
                    elif not (
                        self.classroomSignsOutput
                        or self.dailyScheduleOutput
                        or self.powerpointOutput
                    ):
                        QtWidgets.QMessageBox.warning(
                            None,
                            "No output selected",
                            (
                                "Please select at least one output: "
                                "Classroom Sign, Daily Schedule, or PowerPoint"
                            ),
                        )
                        return
                if os.path.isdir(self.saveReportToPath):
                    result = self.genReportFunction()
                else:
                    QtWidgets.QMessageBox.warning(
                        None,
                        "Save location error",
                        (
                            "The directory you've selected does not exist. "
                            "Please select where you want to save the report to."
                        ),
                    )
                    return
        if result == 0:
            QtWidgets.QMessageBox.warning(
                None,
                "Download error",
                "Report could not be downloaded. Please try again.",
            )
            return
        elif result == 1 and not self.createSigns:
            QtWidgets.QMessageBox.warning(
                None, "Done", "Downloading Destiny Report is complete."
            )
        if self.createSigns:
            chars = set(r"<>?[]:|*")
            if self.saveSignsDirectory == "" or not os.path.isdir(
                self.saveSignsDirectory
            ):
                QtWidgets.QMessageBox.warning(
                    None,
                    "Save location error",
                    "Please select where you want to save the signs to.",
                )
                return
            elif self.genReport:
                if any((c in chars) for c in self.saveSignsDirectory[2:]):
                    print(r"Filename or path contains: <>?[]:|*")
                    QtWidgets.QMessageBox.warning(
                        None,
                        "Save location error",
                        (
                            "The save path cannot contain any of the following "
                            r"characters: <>?[]:| or *"
                        ),
                    )
                    return
                elif not (
                    self.classroomSignsOutput
                    or self.dailyScheduleOutput
                    or self.powerpointOutput
                ):
                    QtWidgets.QMessageBox.warning(
                        None,
                        "No output selected",
                        (
                            "Please select at least one output: "
                            "Classroom Sign, Daily Schedule, or PowerPoint"
                        ),
                    )
                    return
                else:
                    if (
                        self.classroomSignsOutput
                        and self.createSignsFunction(
                            f"{self.saveReportToPath}\\SectionScheduleDailySummary.xls"
                        )
                        == 0
                    ):
                        result = 0
                        pass
                    if (
                        self.dailyScheduleOutput
                        and self.createDailySchedule(
                            f"{self.saveReportToPath}\\SectionScheduleDailySummary.xls"
                        )
                        == 0
                    ):
                        result = 0
                        pass
                    if (
                        self.powerpointOutput
                        and self.createPPT(
                            f"{self.saveReportToPath}\\SectionScheduleDailySummary.xls"
                        )
                        == 0
                    ):
                        result = 0
                        pass

            elif self.useExistingReport and os.path.exists(self.existingReportPath):
                if any((c in chars) for c in self.saveSignsDirectory[2:]):
                    print(r"Filename or path contains: <>?[]:|*")
                    QtWidgets.QMessageBox.warning(
                        None,
                        "Save location error",
                        (
                            "The save path cannot contain any of the following "
                            r"characters: <>?[]:| or *"
                        ),
                    )
                    return
                elif not (
                    self.classroomSignsOutput
                    or self.dailyScheduleOutput
                    or self.powerpointOutput
                ):
                    QtWidgets.QMessageBox.warning(
                        None,
                        "No output selected",
                        (
                            "Please select at least one output: "
                            "Classroom Sign, Daily Schedule, or PowerPoint"
                        ),
                    )
                    return
                else:
                    if (
                        self.classroomSignsOutput
                        and self.createSignsFunction(self.existingReportPath) == 0
                    ):
                        result = 0
                        pass
                    if (
                        self.dailyScheduleOutput
                        and self.createDailySchedule(self.existingReportPath) == 0
                    ):
                        result = 0
                        pass
                    if (
                        self.powerpointOutput
                        and self.createPPT(self.existingReportPath) == 0
                    ):
                        result = 0
                        pass

            elif self.existingReportPath == "" or not os.path.exists(
                self.existingReportPath
            ):
                QtWidgets.QMessageBox.warning(
                    None,
                    "No existing report found!!!",
                    "Please select the location of an existing report.",
                )
                return
            if result == 0:
                QtWidgets.QMessageBox.warning(
                    None, "No classes!!!", "No classes scheduled in date range."
                )
            else:
                QtWidgets.QMessageBox.warning(None, "Done", "Done creating signs.")

    def genReportFunction(self) -> int:
        # Set Chrome defaults to automate download
        service = Service()
        chrome_options = webdriver.ChromeOptions()
        chrome_options.add_experimental_option(
            "prefs",
            {
                "download.default_directory": self.saveReportToPath,
                "download.prompt_for_download": False,
                "download.directory_upgrade": True,
                "safebrowsing.enabled": True,
            },
        )

        # Delete old report if it exists
        if os.path.exists(f"{self.saveReportToPath}\\SectionScheduleDailySummary.xls"):
            os.remove(f"{self.saveReportToPath}\\SectionScheduleDailySummary.xls")

        # Download Destiny Report
        try:
            # browser = webdriver.Chrome(
            #     executable_path=ChromeDriverManager().install(), options=chrome_options
            # )
            browser = webdriver.Chrome(service=service, options=chrome_options)
            browser.get("https://berkeleysv.destinysolutions.com")
            WebDriverWait(browser, 3600).until(
                EC.presence_of_element_located((By.ID, "main-area-body"))
            )
            browser.get(
                "https://berkeleysv.destinysolutions.com/srs/reporting/sectionScheduleDailySummary.do?method=load"  # noqa: E501
            )
            startDateElm = browser.find_element("id", "startDateRecordString")
            startDateElm.send_keys(self.startDate)
            endDateElm = browser.find_element("id", "endDateRecordString")
            endDateElm.send_keys(self.endDate)
            campusElm = browser.find_element("name", "scheduleBlock.campusId")
            campusElm.send_keys(self.center[self.location]["campus"])
            buildingElm = browser.find_element("name", "scheduleBlock.buildingId")
            buildingElm.send_keys(self.center[self.location]["building"])
            outputTypeElm = browser.find_element("name", "outputType")
            outputTypeElm.send_keys("Output to XLS (Export)")
            generateReportElm = browser.find_element("id", "processReport")
            generateReportElm.click()
            while not os.path.exists(
                f"{self.saveReportToPath}\\SectionScheduleDailySummary.xls"
            ):
                time.sleep(1)
            browser.quit()
            return 1
        except WebDriverException:
            browser.quit()
            return 0

    def createSignsFunction(self, reportPath: str) -> int:
        # Read in courses from Excel
        # 1     B   Date
        # 3     D   Type
        # 4     E   Start Time
        # 6     G   End Time
        # 9     J   Section Number
        # 11    L   Section Title
        # 12    M   Instructor
        # 13    N   Building
        # 15    P   Room
        # 16    Q   Configuration
        # 17    R   Technology
        # 18    S   Section Size
        # 20    U   Notes
        # 22    W   Approval Status

        # Read into Pandas dataframe for relevant columns
        pd.set_option("display.max_rows", 500)
        pd.set_option("display.max_columns", 500)
        pd.set_option("display.width", 1000)
        schedule = pd.read_excel(
            reportPath,
            header=6,
            skipfooter=1,
            usecols=[1, 4, 6, 11, 13, 15, 22],
            parse_dates=[1, "Start Time", "End Time"],
            date_format={
                "Date": "%Y/%m/%d %H:%M:%S",
                "Start Time": "%I:%M%p",
                "End Time": "%I:%M%p",
            },
        )
        schedule = schedule[schedule["Approval Status"] == "Final Approval"].copy()

        # Determine if the Destiny report does not have any classes
        if schedule.empty:
            return 0
        # Report is not empty. Determine location and template to use
        else:
            location = self.centerReverse[schedule.iloc[0][4]]["name"]
            template = self.centerReverse[schedule.iloc[0][4]]["template"]
            if location == "SFC":
                self.SFCClassroomSigns(schedule, location, template)
            else:
                self.GBCClassroomSigns(schedule, location, template)
        return 1

    def GBCClassroomSigns(
        self, schedule: pd.DataFrame, location: str, template: str
    ) -> None:
        # Determine the start and end date of the report
        self.startDate = schedule.iloc[0][0].strftime("%Y-%m-%d")
        self.endDate = schedule.iloc[-1][0].strftime("%Y-%m-%d")

        # Sort the raw Destiny Report by Date -> Room # -> Start Time
        schedule["Date"] = schedule["Date"].dt.strftime("%B %d, %Y")
        sortedSchedule = schedule.sort_values(by=["Date", "Room", "Start Time"])
        sortedSchedule["Start Time"] = sortedSchedule["Start Time"].dt.strftime(
            "%I:%M %p"
        )
        sortedSchedule["End Time"] = sortedSchedule["End Time"].dt.strftime("%I:%M %p")

        # Initialize variables for 'for loop', includes formating of the classroom signs
        previousClassroom = ""
        previousDate = ""
        dayofweek = ""
        doc = Document(template)
        paragraph_format = doc.styles["Normal"].paragraph_format
        paragraph_format.space_before = 0
        paragraph_format.space_after = 0
        paragraph_format.line_spacing = 1
        font = doc.styles["Normal"].font
        font.name = "Times New Roman"
        font.bold = True
        font.size = Pt(4)

        # Set page orientation, page size, and margins
        for section in doc.sections:
            section.orientation = WD_ORIENT.LANDSCAPE  # Landscape
            section.page_width = 10058400  # Page Width = 11 inches
            section.page_height = 7772400  # Page Height = 8.5 inches
            section.left_margin = 457200  # Left Margin = 0.5 inches
            section.right_margin = 457200  # Right Margin = 0.5 inches
            section.top_margin = 457200  # Top Margin = 0.5 inches
            section.bottom_margin = 457200  # Bottom Margin = 0.5 inches

        # Add Title, Date, Classroom number, section title and start/end time
        for index in range(0, len(sortedSchedule)):
            newFile = False
            if index != 0 and previousDate != sortedSchedule.iloc[index]["Date"]:
                if index != 0 and index != len(sortedSchedule.index):
                    newFile = True
                    dayofweek = datetime.datetime.strptime(
                        previousDate, "%B %d, %Y"
                    ).strftime("%A")
                    fileDate = datetime.datetime.strptime(
                        previousDate, "%B %d, %Y"
                    ).strftime("%Y-%m-%d")
                    doc.save(
                        f"{self.saveSignsDirectory}\\"
                        f"{location} {fileDate} {dayofweek}.docx"
                    )
                    previousClassroom = ""
                # Create Classroom Signs, set defaults for new file
                doc = Document(template)
                paragraph_format = doc.styles["Normal"].paragraph_format
                paragraph_format.space_before = 0
                paragraph_format.space_after = 0
                paragraph_format.line_spacing = 1
                font = doc.styles["Normal"].font
                font.name = "Times New Roman"
                font.bold = True
                font.size = Pt(4)

                # Set page orientation, page size, and margins
                for section in doc.sections:
                    section.orientation = WD_ORIENT.LANDSCAPE  # Landscape
                    section.page_width = 10058400  # Page Width = 11 inches
                    section.page_height = 7772400  # Page Height = 8.5 inches
                    section.left_margin = 457200  # Left Margin = 0.5 inches
                    section.right_margin = 457200  # Right Margin = 0.5 inches
                    section.top_margin = 457200  # Top Margin = 0.5 inches
                    section.bottom_margin = 457200  # Bottom Margin = 0.5 inches

            if previousClassroom != sortedSchedule.iloc[index]["Room"]:
                if index != 0 and index != len(sortedSchedule.index):
                    if not newFile:
                        doc.add_page_break()  # Reached end of page, start new page

                para = doc.add_paragraph()
                para.alignment = 1
                run = para.add_run(
                    sortedSchedule.iloc[index]["Date"].replace(" 0", " ")
                )  # Date
                run.font.size = Pt(48)

                para = doc.add_paragraph()
                para.alignment = 0
                run = para.add_run(
                    sortedSchedule.iloc[index]["Room"]
                )  # Classroom Number
                run.font.size = Pt(36)

                para = doc.add_paragraph()
                para.alignment = 0
                run = para.add_run("Class:")  # Class
                run.font.size = Pt(36)

                run = para.add_run("\n")
                run.font.size = Pt(2)

                table = doc.add_table(rows=1, cols=2)  # Create table to put each course
                table.alignment = 2
                table.allow_autofit = False

                row = table.rows[0]
                row.cells[0].text = f"{sortedSchedule.iloc[index]['Section Title']}\n"
                row.cells[1].text = (
                    f"{sortedSchedule.iloc[index]['Start Time'].lstrip('0')} to "
                    f"{sortedSchedule.iloc[index]['End Time'].lstrip('0')}"
                )
            else:
                row = table.add_row()  # add a row if course is in same classroom
                row.cells[0].text = f"{sortedSchedule.iloc[index]['Section Title']}\n"
                row.cells[1].text = (
                    f"{sortedSchedule.iloc[index]['Start Time'].lstrip('0')} to "
                    f"{sortedSchedule.iloc[index]['End Time'].lstrip('0')}"
                )

            previousClassroom = sortedSchedule.iloc[index]["Room"]
            previousDate = sortedSchedule.iloc[index]["Date"]

            # Format table columns
            for cell in table.columns[0].cells:
                cell.width = Inches(6.7)
            for cell in table.columns[1].cells:
                cell.width = Inches(3.3)
            # Change font size of text in table
            for row in table.rows:
                for cell in row.cells:
                    paragraphs = cell.paragraphs
                    for paragraph in paragraphs:
                        for run in paragraph.runs:
                            font = run.font
                            font.size = Pt(22)
        # End for loop 'for index in range(0, len(sortedSchedule)):'

        # Save as Microsoft Word docx file for each day
        dayofweek = datetime.datetime.strptime(previousDate, "%B %d, %Y").strftime("%A")
        fileDate = datetime.datetime.strptime(previousDate, "%B %d, %Y").strftime(
            "%Y-%m-%d"
        )
        doc.save(f"{self.saveSignsDirectory}\\{location} {fileDate} {dayofweek}.docx")

    def SFCClassroomSigns(
        self, schedule: pd.DataFrame, location: str, template: str
    ) -> int:
        # Determine the start and end date of the report
        self.startDate = schedule.iloc[0][0].strftime("%Y-%m-%d")
        self.endDate = schedule.iloc[-1][0].strftime("%Y-%m-%d")

        # Sort the raw Destiny Report by Date -> Room # -> Start Time
        schedule["Date"] = schedule["Date"].dt.strftime("%B %d, %Y")
        sortedSchedule = schedule.sort_values(by=["Date", "Room", "Start Time"])
        sortedSchedule["Start Time"] = sortedSchedule["Start Time"].dt.strftime(
            "%I:%M %p"
        )
        sortedSchedule["End Time"] = sortedSchedule["End Time"].dt.strftime("%I:%M %p")

        # Initialize variables for 'for loop', includes formating of the classroom signs
        previousClassroom = ""
        previousDate = ""
        dayofweek = ""
        doc = Document(template)
        doc._body.clear_content()
        paragraph_format = doc.styles["Normal"].paragraph_format
        paragraph_format.space_before = 0
        paragraph_format.space_after = Pt(10)
        font = doc.styles["Normal"].font
        font.name = "Arial"
        font.bold = False
        font.size = Pt(1)

        # Set page orientation, page size, and margins
        for section in doc.sections:
            section.orientation = WD_ORIENT.PORTRAIT  # Landscape
            section.page_width = 7772400  # Page Width = 11 inches
            section.page_height = 10058400  # Page Height = 8.5 inches
            section.left_margin = 457200  # Left Margin = 0.5 inches
            section.right_margin = 457200  # Right Margin = 0.5 inches
            section.top_margin = 457200  # Top Margin = 0.5 inches
            section.bottom_margin = 457200  # Bottom Margin = 0.5 inches

        # Add Title, Date, Classroom number, section title and start/end time
        for index in range(0, len(sortedSchedule)):
            newFile = False
            if index != 0 and previousDate != sortedSchedule.iloc[index]["Date"]:
                if index != 0 and index != len(sortedSchedule.index):
                    newFile = True
                    dayofweek = datetime.datetime.strptime(
                        previousDate, "%B %d, %Y"
                    ).strftime("%A")
                    fileDate = datetime.datetime.strptime(
                        previousDate, "%B %d, %Y"
                    ).strftime("%Y-%m-%d")
                    doc.save(
                        f"{self.saveSignsDirectory}\\"
                        f"{location} {fileDate} {dayofweek}.docx"
                    )
                    previousClassroom = ""
                # Create Classroom Signs, set defaults for new file
                doc = Document(template)
                paragraph_format = doc.styles["Normal"].paragraph_format
                paragraph_format.space_before = 0
                paragraph_format.space_after = Pt(10)
                font = doc.styles["Normal"].font
                font.name = "Arial"
                font.bold = False
                font.size = Pt(1)

                # Set page orientation, page size, and margins
                for section in doc.sections:
                    section.orientation = WD_ORIENT.PORTRAIT  # Landscape
                    section.page_width = 7772400  # Page Width = 11 inches
                    section.page_height = 10058400  # Page Height = 8.5 inches
                    section.left_margin = 457200  # Left Margin = 0.5 inches
                    section.right_margin = 457200  # Right Margin = 0.5 inches
                    section.top_margin = 457200  # Top Margin = 0.5 inches
                    section.bottom_margin = 457200  # Bottom Margin = 0.5 inches

            if previousClassroom != sortedSchedule.iloc[index]["Room"]:
                if index != 0 and index != len(sortedSchedule.index):
                    if not newFile:
                        doc.add_page_break()  # Reached end of page, start new page

                para = doc.add_paragraph()
                para.alignment = 1
                run = para.add_run(
                    datetime.datetime.strptime(
                        sortedSchedule.iloc[index]["Date"], "%B %d, %Y"
                    ).strftime("%A")
                )  # Date
                run.font.size = Pt(30)
                # run.underline = True
                run.bold = True

                para = doc.add_paragraph()
                para.alignment = 1
                run = para.add_run(
                    sortedSchedule.iloc[index]["Room"].replace("Classroom", "Room")
                )  # Classroom Number
                run.font.size = Pt(24)

                run = para.add_run("\n")
                run.font.size = Pt(34)

                table = doc.add_table(rows=1, cols=2)  # Create table to put each course
                table.alignment = 2
                table.allow_autofit = False

                row = table.rows[0]
                row.cells[0].text = "Course"
                row.cells[0].paragraphs[0].runs[0].font.underline = True
                row.cells[0].paragraphs[0].runs[0].font.bold = True
                row.cells[0].paragraphs[0].runs[0].font.size = Pt(22)
                row.cells[1].text = "Time"
                row.cells[1].paragraphs[0].runs[0].font.underline = True
                row.cells[1].paragraphs[0].runs[0].font.bold = True
                row.cells[1].paragraphs[0].runs[0].font.size = Pt(22)

                row = table.add_row()
                row.cells[0].text = f"{sortedSchedule.iloc[index]['Section Title']}\n"
                row.cells[1].text = (
                    f"{sortedSchedule.iloc[index]['Start Time'].lstrip('0')} - "
                    f"{sortedSchedule.iloc[index]['End Time'].lstrip('0')}"
                )

                row.cells[0].paragraphs[0].runs[0].font.size = Pt(22)
                row.cells[1].paragraphs[0].runs[0].font.size = Pt(22)
            else:
                row = table.add_row()  # add a row if course is in same classroom
                row.cells[0].text = f"{sortedSchedule.iloc[index]['Section Title']}\n"
                row.cells[1].text = (
                    f"{sortedSchedule.iloc[index]['Start Time'].lstrip('0')} - "
                    f"{sortedSchedule.iloc[index]['End Time'].lstrip('0')}"
                )
                row.cells[0].paragraphs[0].runs[0].font.size = Pt(22)
                row.cells[1].paragraphs[0].runs[0].font.size = Pt(22)

            previousClassroom = sortedSchedule.iloc[index]["Room"]
            previousDate = sortedSchedule.iloc[index]["Date"]

            # Format table columns
            for cell in table.columns[0].cells:
                cell.width = Inches(5)
            for cell in table.columns[1].cells:
                cell.width = Inches(3.3)

        # Save as Microsoft Word docx file for each day
        dayofweek = datetime.datetime.strptime(previousDate, "%B %d, %Y").strftime("%A")
        fileDate = datetime.datetime.strptime(previousDate, "%B %d, %Y").strftime(
            "%Y-%m-%d"
        )
        doc.save(f"{self.saveSignsDirectory}\\{location} {fileDate} {dayofweek}.docx")
        return 1

    def createDailySchedule(self, reportPath: str) -> int:
        # Read in courses from Excel
        # 1     B   Date
        # 3     D   Type
        # 4     E   Start Time
        # 6     G   End Time
        # 9     J   Section Number
        # 11    L   Section Title
        # 12    M   Instructor
        # 13    N   Building
        # 15    P   Room
        # 16    Q   Configuration
        # 17    R   Technology
        # 18    S   Section Size
        # 20    U   Notes
        # 22    W   Approval Status

        # Read into Pandas dataframe for relevant columns
        schedule = pd.read_excel(
            reportPath,
            header=6,
            skipfooter=1,
            usecols=[1, 4, 6, 9, 11, 12, 13, 15, 22],
            parse_dates=[1, "Start Time", "End Time"],
            date_format={
                "Date": "%Y/%m/%d %H:%M:%S",
                "Start Time": "%I:%M%p",
                "End Time": "%I:%M%p",
            },
        )
        schedule = schedule[schedule["Approval Status"] == "Final Approval"].copy()

        # Determine if the Destiny report does not have any classes
        if schedule.empty:
            return 0
        # Report is not empty. Determine location and template to use
        else:
            location = self.centerReverse[schedule.iloc[0][6]]["name"]
            if location == "SFC":
                self.SFCDailySchedule(schedule, location)
            else:
                self.GBCDailySchedule(schedule, location)
        return 1

    def GBCDailySchedule(self, schedule: pd.DataFrame, location: str) -> int:
        sortedSchedule = schedule.sort_values(by=["Date", "Start Time", "Room"])
        sortedSchedule = sortedSchedule.fillna("")
        dateList = sortedSchedule["Date"].dt.date.unique()

        if len(dateList) == 1:
            writer = pd.ExcelWriter(
                (
                    f"{self.saveSignsDirectory}\\{location} "
                    f"{dateList[0].strftime('%Y-%m-%d')} "
                    f"{dateList[0].strftime('%A')}.xlsx"
                ),
                engine="xlsxwriter",
            )
        else:
            writer = pd.ExcelWriter(
                (
                    f"{self.saveSignsDirectory}\\{location} "
                    f"{dateList[0].strftime('%Y-%m-%d')} "
                    f"{dateList[0].strftime('%A')} to "
                    f"{dateList[-1].strftime('%Y-%m-%d')} "
                    f"{dateList[-1].strftime('%A')}.xlsx"
                ),
                engine="xlsxwriter",
            )
        workbook = writer.book
        # Loop through each day
        for i in range(0, len(dateList)):
            worksheet = workbook.add_worksheet(dateList[i].strftime("%Y-%m-%d"))
            worksheet.set_landscape()  # Page orientation as landscape.
            worksheet.hide_gridlines(0)  # Don’t hide gridlines.
            worksheet.fit_to_pages(1, 1)  # Fit to 1x1 pages.
            worksheet.center_horizontally()
            worksheet.center_vertically()
            worksheet.set_paper(1)  # Set paper size to 8.5" x 11"
            worksheet.set_margins(left=0.25, right=0.25, top=0.25, bottom=0.25)
            worksheet.set_header("", {"margin": 0})
            worksheet.set_footer("", {"margin": 0})

            worksheet.set_column("A:A", 21.5)  # Column A (Start Time) width set to 23.
            worksheet.set_column("B:B", 19)  # Column B (End Time) width set to 19.
            worksheet.set_column("D:D", 64)  # Column D (Section Title) width set to 64.

            titleFormat = workbook.add_format(
                {
                    "font_name": "Verdana",
                    "font_size": 18,
                    "bold": True,
                    "text_wrap": False,
                    "font_color": "#000000",
                }
            )

            headerFormat = workbook.add_format(
                {
                    "font_name": "Verdana",
                    "font_size": 18,
                    "bold": True,
                    "text_wrap": True,
                    "bottom": 2,
                    "bottom_color": "#000000",
                }
            )

            bodyFormat = workbook.add_format(
                {
                    "font_name": "Verdana",
                    "font_size": 18,
                    "bold": False,
                    "valign": "top",
                    "text_wrap": True,
                    "font_color": "#000000",
                }
            )

            worksheet.write(0, 0, "UC Berkeley Extension", titleFormat)
            worksheet.write(
                0,
                4,
                (
                    f"{dateList[i].strftime('%A')} "
                    f"{dateList[i].strftime('%B %d, %Y').replace(' 0', ' ')}"
                ),
                titleFormat,
            )
            for col_num, value in enumerate(
                [
                    "Start Time",
                    "End Time",
                    "Section Number",
                    "Section Title",
                    "Instructor",
                    "Room",
                ]
            ):
                worksheet.write(2, col_num, value, headerFormat)

            singleDaySched = sortedSchedule.loc[
                sortedSchedule["Date"].dt.date == dateList[i], :
            ]

            morningBlock = (
                singleDaySched.set_index("Start Time")
                .between_time("00:00", "12:00", inclusive="left")
                .reset_index()
            )
            afternoonBlock = (
                singleDaySched.set_index("Start Time")
                .between_time("12:00", "17:00", inclusive="left")
                .reset_index()
            )
            eveningBlock = (
                singleDaySched.set_index("Start Time")
                .between_time("17:00", "00:00", inclusive="left")
                .reset_index()
            )

            morningBlock["Start Time"] = morningBlock["Start Time"].dt.strftime(
                "%I:%M %p"
            )
            morningBlock["End Time"] = morningBlock["End Time"].dt.strftime("%I:%M %p")
            afternoonBlock["Start Time"] = afternoonBlock["Start Time"].dt.strftime(
                "%I:%M %p"
            )
            afternoonBlock["End Time"] = afternoonBlock["End Time"].dt.strftime(
                "%I:%M %p"
            )
            eveningBlock["Start Time"] = eveningBlock["Start Time"].dt.strftime(
                "%I:%M %p"
            )
            eveningBlock["End Time"] = eveningBlock["End Time"].dt.strftime("%I:%M %p")

            excelRow = 3
            if not morningBlock.empty:
                worksheet.write(excelRow, 0, "Morning Classes", titleFormat)
                excelRow += 1
                for i, row in morningBlock.iterrows():
                    worksheet.write(
                        excelRow, 0, row["Start Time"].lstrip("0"), bodyFormat
                    )
                    worksheet.write(
                        excelRow, 1, row["End Time"].lstrip("0"), bodyFormat
                    )
                    worksheet.write(excelRow, 2, row["Section Number"], bodyFormat)
                    worksheet.write(excelRow, 3, row["Section Title"], bodyFormat)
                    if row["Instructor"] == "Instructor To Be Announced":
                        worksheet.write(excelRow, 4, "TBA", bodyFormat)
                    elif not pd.isnull(row["Instructor"]):
                        worksheet.write(excelRow, 4, row["Instructor"], bodyFormat)
                    worksheet.write(excelRow, 5, row["Room"], bodyFormat)
                    excelRow += 1

            if not afternoonBlock.empty:
                excelRow += 1
                worksheet.write(excelRow, 0, "Afternoon Classes", titleFormat)
                excelRow += 1
                for i, row in afternoonBlock.iterrows():
                    worksheet.write(
                        excelRow, 0, row["Start Time"].lstrip("0"), bodyFormat
                    )
                    worksheet.write(
                        excelRow, 1, row["End Time"].lstrip("0"), bodyFormat
                    )
                    worksheet.write(excelRow, 2, row["Section Number"], bodyFormat)
                    worksheet.write(excelRow, 3, row["Section Title"], bodyFormat)
                    if row["Instructor"] == "Instructor To Be Announced":
                        worksheet.write(excelRow, 4, "TBA", bodyFormat)
                    elif not pd.isnull(row["Instructor"]):
                        worksheet.write(excelRow, 4, row["Instructor"], bodyFormat)
                    worksheet.write(excelRow, 5, row["Room"], bodyFormat)
                    excelRow += 1

            if not eveningBlock.empty:
                excelRow += 1
                worksheet.write(excelRow, 0, "Evening Classes", titleFormat)
                excelRow += 1
                for i, row in eveningBlock.iterrows():
                    worksheet.write(
                        excelRow, 0, row["Start Time"].lstrip("0"), bodyFormat
                    )
                    worksheet.write(
                        excelRow, 1, row["End Time"].lstrip("0"), bodyFormat
                    )
                    worksheet.write(excelRow, 2, row["Section Number"], bodyFormat)
                    worksheet.write(excelRow, 3, row["Section Title"], bodyFormat)
                    if row["Instructor"] == "Instructor To Be Announced":
                        worksheet.write(excelRow, 4, "TBA", bodyFormat)
                    elif not pd.isnull(row["Instructor"]):
                        worksheet.write(excelRow, 4, row["Instructor"], bodyFormat)
                    worksheet.write(excelRow, 5, row["Room"], bodyFormat)
                    excelRow += 1

            # Set the minimum column width of ['Section Number', 'Instructor', 'Room']
            max_length = [
                14,
                10,
                4,
            ]
            for i, row in singleDaySched.iterrows():
                try:  # Necessary to avoid error on empty cells
                    if len(str(row["Section Number"])) > max_length[0]:
                        max_length[0] = len(row["Section Number"])
                    if (
                        (row["Instructor"] != "Instructor To Be Announced")
                        and (len(str(row["Instructor"])) > max_length[1])
                        and (len(str(row["Instructor"])) <= 20)
                    ):
                        max_length[1] = len(row["Instructor"])
                    if len(str(row["Room"])) > max_length[2]:
                        max_length[2] = len(row["Room"])
                except Exception:
                    pass

            # Auto fit columns based on character count.
            worksheet.set_column("C:C", (max_length[0] + 1) * (40 / 19))
            worksheet.set_column("E:E", (max_length[1] + (8 / 9)) * (27 / 14))
            worksheet.set_column("F:F", (max_length[2] - (7 / 13)) * (13 / 6))

        workbook.close()
        return 1

    def SFCDailySchedule(self, schedule: pd.DataFrame, location: str) -> int:
        sortedSchedule = schedule.sort_values(by=["Date", "Start Time", "Room"])
        sortedSchedule = sortedSchedule.fillna("")
        dateList = sortedSchedule["Date"].dt.date.unique()

        if len(dateList) == 1:
            writer = pd.ExcelWriter(
                (
                    f"{self.saveSignsDirectory}\\{location} "
                    f"{dateList[0].strftime('%Y-%m-%d')} "
                    f"{dateList[0].strftime('%A')}.xlsx"
                ),
                engine="xlsxwriter",
            )
        else:
            writer = pd.ExcelWriter(
                (
                    f"{self.saveSignsDirectory}\\{location} "
                    f"{dateList[0].strftime('%Y-%m-%d')} "
                    f"{dateList[0].strftime('%A')} to "
                    f"{dateList[-1].strftime('%Y-%m-%d')} "
                    f"{dateList[-1].strftime('%A')}.xlsx"
                ),
                engine="xlsxwriter",
            )
        workbook = writer.book
        # Loop through each day
        for i in range(0, len(dateList)):
            worksheet = workbook.add_worksheet(dateList[i].strftime("%Y-%m-%d"))
            worksheet.set_portrait()  # Page orientation as landscape.
            # worksheet.hide_gridlines(0)     # Don’t hide gridlines.
            worksheet.fit_to_pages(1, 1)  # Fit to 1x1 pages.
            worksheet.center_horizontally()
            worksheet.center_vertically()
            worksheet.set_paper(1)  # Set paper size to 8.5" x 11"
            worksheet.set_margins(left=0.25, right=0.25, top=0.25, bottom=0.25)
            worksheet.set_header("", {"margin": 0})
            worksheet.set_footer("", {"margin": 0})

            worksheet.set_column("A:A", 21.5)  # Column A (Start Time) width set to 23.
            worksheet.set_column("B:B", 19)  # Column B (End Time) width set to 19.
            worksheet.set_column("D:D", 64)  # Column D (Section Title) width set to 64.

            titleFormat = workbook.add_format(
                {
                    "font_name": "Arial",
                    "font_size": 30,
                    "align": "center",
                    "bold": True,
                    "text_wrap": False,
                    "font_color": "red",
                    "underline": True,
                }
            )

            blankFormat = workbook.add_format(
                {
                    "font_name": "Arial",
                    "font_size": 8,
                }
            )

            headerFormat = workbook.add_format(
                {
                    "font_name": "Arial",
                    "font_size": 24,
                    "bold": True,
                    "text_wrap": True,
                    "underline": True,
                }
            )

            blockFormat = workbook.add_format(
                {
                    "font_name": "Arial",
                    "font_size": 24,
                    "bold": True,
                    "text_wrap": False,
                    "align": "left",
                    "valign": "top",
                    "underline": True,
                }
            )

            bodyFormat = workbook.add_format(
                {
                    "font_name": "Arial",
                    "font_size": 21,
                    "bold": False,
                    "valign": "top",
                    "text_wrap": True,
                    "font_color": "#000000",
                }
            )

            # Title and header rows
            worksheet.merge_range(
                "A1:E1",
                (
                    f"UC Berkeley Extension - {dateList[i].strftime('%A')}, "
                    f"{dateList[i].strftime('%B %d, %Y').replace(' 0', ' ')}"
                ),
                titleFormat,
            )
            for col_num, value in enumerate(
                ["Start Time", "End Time", "Section Number", "Section Title", "Room"]
            ):
                worksheet.write(1, col_num, value, headerFormat)

            # Filter by daytime and evening, then by floor
            singleDaySched = sortedSchedule.loc[
                sortedSchedule["Date"].dt.date == dateList[i], :
            ]
            daytimeBlock = (
                singleDaySched.set_index("Start Time")
                .between_time("00:00", "17:00", inclusive="left")
                .reset_index()
            )
            daytimeBlock = daytimeBlock.sort_values(by=["Room", "Start Time"])
            daytimeBlock["Start Time"] = daytimeBlock["Start Time"].dt.strftime(
                "%I:%M %p"
            )
            daytimeBlock["End Time"] = daytimeBlock["End Time"].dt.strftime("%I:%M %p")
            daytime5thFlr = daytimeBlock.loc[daytimeBlock["Room"] <= "Classroom 515", :]
            daytime6thFlr = daytimeBlock.loc[
                (daytimeBlock["Room"] >= "Classroom 602")
                & (daytimeBlock["Room"] <= "Classroom 613"),
                :,
            ]
            daytime7thFlr = daytimeBlock.loc[daytimeBlock["Room"] >= "Classroom 702", :]

            eveningBlock = (
                singleDaySched.set_index("Start Time")
                .between_time("17:00", "00:00", inclusive="left")
                .reset_index()
            )
            eveningBlock["Start Time"] = eveningBlock["Start Time"].dt.strftime(
                "%I:%M %p"
            )
            eveningBlock["End Time"] = eveningBlock["End Time"].dt.strftime("%I:%M %p")
            eveningBlock = eveningBlock.sort_values(by=["Room", "Start Time"])
            evening5thFlr = eveningBlock.loc[eveningBlock["Room"] <= "Classroom 515", :]
            evening6thFlr = eveningBlock.loc[
                (eveningBlock["Room"] >= "Classroom 602")
                & (eveningBlock["Room"] <= "Classroom 613"),
                :,
            ]
            evening7thFlr = eveningBlock.loc[eveningBlock["Room"] >= "Classroom 702", :]

            # Write to cells starting with daytime courses, by floor
            excelRow = 2
            worksheet.merge_range(
                f"A{str(excelRow+1)}:E{str(excelRow +1)}", "", blankFormat
            )
            if not daytimeBlock.empty:
                excelRow += 1
                worksheet.merge_range(
                    f"A{str(excelRow + 1)}:B{str(excelRow + 2)}",
                    "Daytime Classes",
                    blockFormat,
                )
                worksheet.merge_range(
                    f"C{str(excelRow + 1)}:E{str(excelRow + 2)}", "", blockFormat
                )
                excelRow += 1

                if not daytime5thFlr.empty:
                    excelRow += 1
                    worksheet.merge_range(
                        f"A{str(excelRow+1)}:B{str(excelRow +1)}",
                        "5th Floor",
                        blockFormat,
                    )
                    worksheet.merge_range(
                        f"C{str(excelRow + 1)}:E{str(excelRow + 1)}", "", blockFormat
                    )
                    excelRow += 1

                    for i, row in daytime5thFlr.iterrows():
                        worksheet.write(
                            excelRow, 0, row["Start Time"].lstrip("0"), bodyFormat
                        )
                        worksheet.write(
                            excelRow, 1, row["End Time"].lstrip("0"), bodyFormat
                        )
                        worksheet.write(excelRow, 2, row["Section Number"], bodyFormat)
                        worksheet.write(excelRow, 3, row["Section Title"], bodyFormat)
                        worksheet.write(
                            excelRow,
                            4,
                            row["Room"].replace("Classroom", ""),
                            bodyFormat,
                        )
                        excelRow += 1
                    worksheet.merge_range(
                        f"A{str(excelRow+1)}:E{str(excelRow +1)}", "", blankFormat
                    )

                if not daytime6thFlr.empty:
                    excelRow += 1
                    worksheet.merge_range(
                        f"A{str(excelRow+1)}:B{str(excelRow +1)}",
                        "6th Floor",
                        blockFormat,
                    )
                    worksheet.merge_range(
                        f"C{str(excelRow + 1)}:E{str(excelRow + 1)}", "", blockFormat
                    )
                    excelRow += 1

                    for i, row in daytime6thFlr.iterrows():
                        worksheet.write(
                            excelRow, 0, row["Start Time"].lstrip("0"), bodyFormat
                        )
                        worksheet.write(
                            excelRow, 1, row["End Time"].lstrip("0"), bodyFormat
                        )
                        worksheet.write(excelRow, 2, row["Section Number"], bodyFormat)
                        worksheet.write(excelRow, 3, row["Section Title"], bodyFormat)
                        worksheet.write(
                            excelRow,
                            4,
                            row["Room"].replace("Classroom", ""),
                            bodyFormat,
                        )
                        excelRow += 1
                    worksheet.merge_range(
                        f"A{str(excelRow+1)}:E{str(excelRow +1)}", "", blankFormat
                    )

                if not daytime7thFlr.empty:
                    excelRow += 1
                    worksheet.merge_range(
                        f"A{str(excelRow+1)}:B{str(excelRow +1)}",
                        "7th Floor",
                        blockFormat,
                    )
                    worksheet.merge_range(
                        f"C{str(excelRow + 1)}:E{str(excelRow + 1)}", "", blockFormat
                    )
                    excelRow += 1

                    for i, row in daytime7thFlr.iterrows():
                        worksheet.write(
                            excelRow, 0, row["Start Time"].lstrip("0"), bodyFormat
                        )
                        worksheet.write(
                            excelRow, 1, row["End Time"].lstrip("0"), bodyFormat
                        )
                        worksheet.write(excelRow, 2, row["Section Number"], bodyFormat)
                        worksheet.write(excelRow, 3, row["Section Title"], bodyFormat)
                        worksheet.write(
                            excelRow,
                            4,
                            row["Room"].replace("Classroom", ""),
                            bodyFormat,
                        )
                        excelRow += 1
                    worksheet.merge_range(
                        f"A{str(excelRow+1)}:E{str(excelRow +1)}", "", blankFormat
                    )

            # Write to cells for evening courses, by floor
            if not eveningBlock.empty:
                excelRow += 1
                worksheet.merge_range(
                    f"A{str(excelRow+1)}:B{str(excelRow + 2)}",
                    "Evening Classes",
                    blockFormat,
                )
                worksheet.merge_range(
                    f"C{str(excelRow + 1)}:E{str(excelRow + 2)}", "", blockFormat
                )
                excelRow += 1

                if not evening5thFlr.empty:
                    excelRow += 1
                    worksheet.merge_range(
                        f"A{str(excelRow+1)}:B{str(excelRow +1)}",
                        "5th Floor",
                        blockFormat,
                    )
                    worksheet.merge_range(
                        f"C{str(excelRow + 1)}:E{str(excelRow + 1)}", "", blockFormat
                    )
                    excelRow += 1

                    for i, row in evening5thFlr.iterrows():
                        worksheet.write(
                            excelRow, 0, row["Start Time"].lstrip("0"), bodyFormat
                        )
                        worksheet.write(
                            excelRow, 1, row["End Time"].lstrip("0"), bodyFormat
                        )
                        worksheet.write(excelRow, 2, row["Section Number"], bodyFormat)
                        worksheet.write(excelRow, 3, row["Section Title"], bodyFormat)
                        worksheet.write(
                            excelRow,
                            4,
                            row["Room"].replace("Classroom", ""),
                            bodyFormat,
                        )
                        excelRow += 1
                    worksheet.merge_range(
                        f"A{str(excelRow+1)}:E{str(excelRow +1)}", "", blankFormat
                    )

                if not evening6thFlr.empty:
                    excelRow += 1
                    worksheet.merge_range(
                        f"A{str(excelRow+1)}:B{str(excelRow +1)}",
                        "6th Floor",
                        blockFormat,
                    )
                    worksheet.merge_range(
                        f"C{str(excelRow + 1)}:E{str(excelRow + 1)}", "", blockFormat
                    )
                    excelRow += 1

                    for i, row in evening6thFlr.iterrows():
                        worksheet.write(
                            excelRow, 0, row["Start Time"].lstrip("0"), bodyFormat
                        )
                        worksheet.write(
                            excelRow, 1, row["End Time"].lstrip("0"), bodyFormat
                        )
                        worksheet.write(excelRow, 2, row["Section Number"], bodyFormat)
                        worksheet.write(excelRow, 3, row["Section Title"], bodyFormat)
                        worksheet.write(
                            excelRow,
                            4,
                            row["Room"].replace("Classroom", ""),
                            bodyFormat,
                        )
                        excelRow += 1
                    worksheet.merge_range(
                        f"A{str(excelRow+1)}:E{str(excelRow +1)}", "", blankFormat
                    )

                if not evening7thFlr.empty:
                    excelRow += 1
                    worksheet.merge_range(
                        f"A{str(excelRow+1)}:B{str(excelRow +1)}",
                        "7th Floor",
                        blockFormat,
                    )
                    worksheet.merge_range(
                        f"C{str(excelRow + 1)}:E{str(excelRow + 1)}", "", blockFormat
                    )
                    excelRow += 1

                    for i, row in evening7thFlr.iterrows():
                        worksheet.write(
                            excelRow, 0, row["Start Time"].lstrip("0"), bodyFormat
                        )
                        worksheet.write(
                            excelRow, 1, row["End Time"].lstrip("0"), bodyFormat
                        )
                        worksheet.write(excelRow, 2, row["Section Number"], bodyFormat)
                        worksheet.write(excelRow, 3, row["Section Title"], bodyFormat)
                        worksheet.write(
                            excelRow,
                            4,
                            row["Room"].replace("Classroom", ""),
                            bodyFormat,
                        )
                        excelRow += 1

            # Adjust column with of the excel file
            worksheet.set_column("A:A", 18.57)
            worksheet.set_column("B:B", 18.57)
            worksheet.set_column("C:C", 43)
            worksheet.set_column("D:D", 100)
            worksheet.set_column("E:E", 13.86)

        workbook.close()
        return 1

    def createPPT(self, reportPath: str) -> int:
        # Read in courses from Excel
        # 1     B   Date
        # 3     D   Type
        # 4     E   Start Time
        # 6     G   End Time
        # 9     J   Section Number
        # 11    L   Section Title
        # 12    M   Instructor
        # 13    N   Building
        # 15    P   Room
        # 16    Q   Configuration
        # 17    R   Technology
        # 18    S   Section Size
        # 20    U   Notes
        # 22    W   Approval Status

        # Read into Pandas dataframe for relevant columns
        schedule = pd.read_excel(
            reportPath,
            header=6,
            skipfooter=1,
            usecols=[1, 4, 6, 9, 11, 12, 13, 15, 22],
            parse_dates=[1, "Start Time", "End Time"],
            date_format={
                "Date": "%Y/%m/%d %H:%M:%S",
                "Start Time": "%I:%M%p",
                "End Time": "%I:%M%p",
            },
        )
        schedule = schedule[schedule["Approval Status"] == "Final Approval"].copy()

        # Determine if the Destiny report does not have any classes
        if schedule.empty:
            return 0
        # Report is not empty. Determine location and which template to use
        else:
            location = self.centerReverse[schedule.iloc[0][6]]["name"]
            template = self.centerReverse[schedule.iloc[0][6]]["pptTemplate"]
            if location == "SFC":
                self.SFCppt(schedule, location, template)
            else:
                self.GBCppt(schedule, location, template)
        return 1

    def GBCppt(self, schedule: pd.DataFrame, location: str, template: str) -> int:
        # Sort the schedule
        sortedSchedule = schedule.sort_values(by=["Date", "Start Time", "Room"])
        sortedSchedule = sortedSchedule.fillna("")
        dateList = sortedSchedule["Date"].dt.date.unique()

        # Upload GBC schedule if setting and URL are set in config.ini file
        if self.uploadGBCSchedule and self.GBCScheduleURL:
            self.GBCScheduleToGSheets(
                dateList[0],
                sortedSchedule.loc[sortedSchedule["Date"].dt.date == dateList[0], :],
            )

        # Write out schedule one block per slide. Hide slide if no classes.
        for i in range(0, len(dateList)):
            singleDaySched = sortedSchedule.loc[
                sortedSchedule["Date"].dt.date == dateList[i], :
            ]

            morningBlock = (
                singleDaySched.set_index("Start Time")
                .between_time("00:00", "12:00", inclusive="left")
                .reset_index()
            )
            afternoonBlock = (
                singleDaySched.set_index("Start Time")
                .between_time("12:00", "17:00", inclusive="left")
                .reset_index()
            )
            eveningBlock = (
                singleDaySched.set_index("Start Time")
                .between_time("17:00", "00:00", inclusive="left")
                .reset_index()
            )

            morningBlock["Start Time"] = morningBlock["Start Time"].dt.strftime(
                "%I:%M %p"
            )
            morningBlock["End Time"] = morningBlock["End Time"].dt.strftime("%I:%M %p")
            afternoonBlock["Start Time"] = afternoonBlock["Start Time"].dt.strftime(
                "%I:%M %p"
            )
            afternoonBlock["End Time"] = afternoonBlock["End Time"].dt.strftime(
                "%I:%M %p"
            )
            eveningBlock["Start Time"] = eveningBlock["Start Time"].dt.strftime(
                "%I:%M %p"
            )
            eveningBlock["End Time"] = eveningBlock["End Time"].dt.strftime("%I:%M %p")

            prs = Presentation(template)

            slide = prs.slides[1]
            if not morningBlock.empty:
                text_frame = slide.shapes[2].text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = (
                    f"UC Berkeley Extension {dateList[i].strftime('%A')}, "
                    f"{dateList[i].strftime('%B %d, %Y').replace(' 0', ' ')}"
                )
                font = run.font
                font.size = Pt(120)
                font.name = "Calibri"
                font.bold = True
                font.color.rgb = RGBColor(0xFF, 0xFF, 0x00)
                table = slide.shapes[0].table

                currentRow = 1
                colorToggle = True
                for idx, row in morningBlock.iterrows():
                    if colorToggle:
                        color = RGBColor(0xFF, 0xFF, 0xFF)
                    else:
                        color = RGBColor(0xFA, 0xC0, 0x90)
                    colorToggle = not colorToggle
                    for col, value in enumerate(
                        [
                            "Start Time",
                            "End Time",
                            "Section Title",
                            "Instructor",
                            "Room",
                        ]
                    ):
                        text_frame = table.rows[currentRow].cells[col].text_frame
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        if (value == "Instructor") and (
                            (row["Instructor"] == "Instructor To Be Announced")
                            or (pd.isnull(row["Instructor"]))
                        ):
                            run.text = "TBA"
                        else:
                            run.text = str(row[value]).lstrip("0")
                        font = run.font
                        fontSize = (
                            65
                            if 950 // len(morningBlock.index) >= 65
                            else 950 // len(morningBlock.index)
                        )
                        font.size = Pt(fontSize)
                        font.name = "Calibri"
                        font.bold = True
                        font.color.rgb = color
                    currentRow += 1
            else:
                slide._element.set("show", "0")

            slide = prs.slides[2]
            if not afternoonBlock.empty:
                text_frame = slide.shapes[2].text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = (
                    f"UC Berkeley Extension {dateList[i].strftime('%A')} "
                    f"{dateList[i].strftime('%B %d, %Y').replace(' 0', ' ')}"
                )
                font = run.font
                font.size = Pt(120)
                font.name = "Calibri"
                font.bold = True
                font.color.rgb = RGBColor(0xFF, 0xFF, 0x00)
                table = slide.shapes[0].table

                currentRow = 1
                colorToggle = True
                for idx, row in afternoonBlock.iterrows():
                    if colorToggle:
                        color = RGBColor(0xFF, 0xFF, 0xFF)
                    else:
                        color = RGBColor(0xFA, 0xC0, 0x90)
                    colorToggle = not colorToggle
                    for col, value in enumerate(
                        [
                            "Start Time",
                            "End Time",
                            "Section Title",
                            "Instructor",
                            "Room",
                        ]
                    ):
                        text_frame = table.rows[currentRow].cells[col].text_frame
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        if (value == "Instructor") and (
                            (row["Instructor"] == "Instructor To Be Announced")
                            or (pd.isnull(row["Instructor"]))
                        ):
                            run.text = "TBA"
                        else:
                            run.text = str(row[value]).lstrip("0")
                        font = run.font
                        fontSize = (
                            65
                            if 950 // len(afternoonBlock.index) >= 65
                            else 950 // len(afternoonBlock.index)
                        )
                        font.size = Pt(fontSize)
                        font.name = "Calibri"
                        font.bold = True
                        font.color.rgb = color
                    currentRow += 1
            else:
                slide._element.set("show", "0")

            slide = prs.slides[3]
            if not eveningBlock.empty:
                text_frame = slide.shapes[2].text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = (
                    f"UC Berkeley Extension {dateList[i].strftime('%A')} "
                    f"{dateList[i].strftime('%B %d, %Y').replace(' 0', ' ')}"
                )
                font = run.font
                font.size = Pt(120)
                font.name = "Calibri"
                font.bold = True
                font.color.rgb = RGBColor(0xFF, 0xFF, 0x00)
                table = slide.shapes[0].table

                currentRow = 1
                colorToggle = True
                for idx, row in eveningBlock.iterrows():
                    if colorToggle:
                        color = RGBColor(0xFF, 0xFF, 0xFF)
                    else:
                        color = RGBColor(0xFA, 0xC0, 0x90)
                    colorToggle = not colorToggle
                    for col, value in enumerate(
                        [
                            "Start Time",
                            "End Time",
                            "Section Title",
                            "Instructor",
                            "Room",
                        ]
                    ):
                        text_frame = table.rows[currentRow].cells[col].text_frame
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        if (value == "Instructor") and (
                            (row["Instructor"] == "Instructor To Be Announced")
                            or (pd.isnull(row["Instructor"]))
                        ):
                            run.text = "TBA"
                        else:
                            run.text = str(row[value]).lstrip("0")
                        font = run.font
                        fontSize = (
                            65
                            if 950 // len(eveningBlock.index) >= 65
                            else 950 // len(eveningBlock.index)
                        )
                        font.size = Pt(fontSize)
                        font.name = "Calibri"
                        font.bold = True
                        font.color.rgb = color
                    currentRow += 1
            else:
                slide._element.set("show", "0")

            # Save as Microsoft Powerpoint pptx file, per day
            prs.save(
                (
                    f"{self.saveSignsDirectory}\\{location} "
                    f"{dateList[i].strftime('%Y-%m-%d')} "
                    f"{dateList[i].strftime('%A')}.pptx"
                )
            )

        return 1

    def GBCScheduleToGSheets(
        self, date: datetime.datetime, schedule: pd.DataFrame
    ) -> pygsheets.PyGsheetsException:
        # Sort the schedule by time of day blocks
        morningBlock = (
            schedule.set_index("Start Time")
            .between_time("00:00", "12:00", inclusive="left")
            .reset_index()
        )
        afternoonBlock = (
            schedule.set_index("Start Time")
            .between_time("12:00", "17:00", inclusive="left")
            .reset_index()
        )
        eveningBlock = (
            schedule.set_index("Start Time")
            .between_time("17:00", "00:00", inclusive="left")
            .reset_index()
        )

        morningBlock["Start Time"] = morningBlock["Start Time"].dt.strftime("%I:%M %p")
        morningBlock["End Time"] = morningBlock["End Time"].dt.strftime("%I:%M %p")
        afternoonBlock["Start Time"] = afternoonBlock["Start Time"].dt.strftime(
            "%I:%M %p"
        )
        afternoonBlock["End Time"] = afternoonBlock["End Time"].dt.strftime("%I:%M %p")
        eveningBlock["Start Time"] = eveningBlock["Start Time"].dt.strftime("%I:%M %p")
        eveningBlock["End Time"] = eveningBlock["End Time"].dt.strftime("%I:%M %p")

        blockList = [
            ("Morning", morningBlock),
            ("Afternoon", afternoonBlock),
            ("Evening", eveningBlock),
        ]
        dirpath = os.getcwd()

        # Connect to Google Sheets and update with current schedule
        client = pygsheets.authorize(service_file=f"{dirpath}/service_file.json")
        try:
            sheet = client.open_by_url(self.GBCScheduleURL)
        except pygsheets.SpreadsheetNotFound as error:
            print(error)
            return error
        finally:
            for i in range(len(blockList)):
                wks = sheet.worksheet_by_title(blockList[i][0])
                wks.clear(start="A1", end=None, fields="*")
                wks.update_value(
                    "A1",
                    (
                        f"UC Berkeley Extension - {date.strftime('%A')}, "
                        f"{date.strftime('%B %d, %Y').replace(' 0', ' ')}"
                    ),
                )
                wks.set_dataframe(
                    blockList[i][1][
                        [
                            "Start Time",
                            "End Time",
                            "Section Title",
                            "Instructor",
                            "Room",
                        ]
                    ],
                    "A2",
                    fit=True,
                )

    def SFCppt(self, schedule: pd.DataFrame, location: str, template: str) -> int:
        # Sort the schedule
        sortedSchedule = schedule.sort_values(by=["Date", "Start Time", "Room"])
        sortedSchedule = sortedSchedule.fillna("")
        dateList = sortedSchedule["Date"].dt.date.unique()

        # Upload SFC schedule if setting and URL are set in config.ini file
        if self.uploadSFCSchedule and self.SFCScheduleURL:
            self.SFCScheduleToGSheets(
                dateList[0],
                sortedSchedule.loc[sortedSchedule["Date"].dt.date == dateList[0], :],
            )

        # Write out schedule one block per slide. Hide slide if no classes.
        for i in range(0, len(dateList)):
            singleDaySched = sortedSchedule.loc[
                sortedSchedule["Date"].dt.date == dateList[i], :
            ]
            daytimeBlock = (
                singleDaySched.set_index("Start Time")
                .between_time("00:00", "17:00", inclusive="left")
                .reset_index()
            )
            daytimeBlock["Start Time"] = daytimeBlock["Start Time"].dt.strftime(
                "%I:%M %p"
            )
            daytimeBlock["End Time"] = daytimeBlock["End Time"].dt.strftime("%I:%M %p")
            daytimeBlock = daytimeBlock.sort_values(by=["Room", "Start Time"])
            daytime5thFlr = daytimeBlock.loc[daytimeBlock["Room"] <= "Classroom 515", :]
            daytime6thFlr = daytimeBlock.loc[
                (daytimeBlock["Room"] >= "Classroom 602")
                & (daytimeBlock["Room"] <= "Classroom 613"),
                :,
            ]
            daytime7thFlr = daytimeBlock.loc[daytimeBlock["Room"] >= "Classroom 702", :]

            eveningBlock = (
                singleDaySched.set_index("Start Time")
                .between_time("17:00", "00:00", inclusive="left")
                .reset_index()
            )
            eveningBlock["Start Time"] = eveningBlock["Start Time"].dt.strftime(
                "%I:%M %p"
            )
            eveningBlock["End Time"] = eveningBlock["End Time"].dt.strftime("%I:%M %p")
            eveningBlock = eveningBlock.sort_values(by=["Room", "Start Time"])
            evening5thFlr = eveningBlock.loc[eveningBlock["Room"] <= "Classroom 515", :]
            evening6thFlr = eveningBlock.loc[
                (eveningBlock["Room"] >= "Classroom 602")
                & (eveningBlock["Room"] <= "Classroom 613"),
                :,
            ]
            evening7thFlr = eveningBlock.loc[eveningBlock["Room"] >= "Classroom 702", :]

            # Max font size Pt(60). Scale font size down based on number of rows used.
            daytimeRowCount = (
                len(daytime5thFlr.index)
                + 2 * (not daytime5thFlr.empty)
                + len(daytime6thFlr.index)
                + 2 * (not daytime6thFlr.empty)
                + len(daytime7thFlr.index)
                + 1 * (not daytime7thFlr.empty)
            )
            daytimeFontSize = -1.0603 * daytimeRowCount + 72.336
            if daytimeFontSize > 60:
                daytimeFontSize = 60
            eveningRowCount = (
                len(evening5thFlr.index)
                + 2 * (not evening5thFlr.empty)
                + len(evening6thFlr.index)
                + 2 * (not evening6thFlr.empty)
                + len(evening7thFlr.index)
                + 1 * (not evening7thFlr.empty)
            )
            eveningFontSize = -1.0603 * eveningRowCount + 72.336
            if eveningFontSize > 60:
                eveningFontSize = 60

            prs = Presentation(template)

            slide = prs.slides[1]
            if not daytimeBlock.empty:
                text_frame = slide.shapes[2].text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = (
                    f"UC Berkeley Extension - {dateList[i].strftime('%A')}, "
                    f"{dateList[i].strftime('%B %d, %Y').replace(' 0', ' ')}"
                )
                font = run.font
                font.size = Pt(70)
                font.name = "Calibri"
                font.bold = True
                font.color.rgb = RGBColor(0xFF, 0xFF, 0x00)
                table = slide.shapes[0].table

                currentRow = 0
                if not daytime5thFlr.empty:
                    text_frame = table.rows[currentRow].cells[0].text_frame
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = "5th Floor"
                    font = run.font
                    font.size = Pt(int(daytimeFontSize))
                    font.name = "Arial"
                    font.bold = True
                    font.underline = True
                    font.color.rgb = RGBColor(0xFF, 0xFF, 0x00)

                    text_frame = table.rows[currentRow].cells[4].text_frame
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = "Room"
                    font = run.font
                    font.size = Pt(int(daytimeFontSize))
                    font.name = "Arial"
                    font.bold = True
                    font.underline = True
                    font.color.rgb = RGBColor(0xFF, 0xFF, 0x00)

                    currentRow += 1
                    colorToggle = True

                    for idx, row in daytime5thFlr.iterrows():
                        if colorToggle:
                            color = RGBColor(0xFF, 0xFF, 0xFF)
                        else:
                            color = RGBColor(0xFA, 0xC0, 0x90)
                        colorToggle = not colorToggle
                        for col, value in enumerate(
                            [
                                "Start Time",
                                "End Time",
                                "Section Number",
                                "Section Title",
                                "Room",
                            ]
                        ):
                            text_frame = table.rows[currentRow].cells[col].text_frame
                            p = text_frame.paragraphs[0]
                            run = p.add_run()
                            if value == "Start Time" or value == "End Time":
                                run.text = str(row[value]).lstrip("0")
                            elif value == "Room":
                                run.text = str(row[value]).lstrip("Classroom")
                            else:
                                run.text = str(row[value])
                            font = run.font
                            font.size = Pt(int(daytimeFontSize))
                            font.name = "Arial"
                            font.bold = True
                            font.color.rgb = color
                        currentRow += 1
                    currentRow += 1

                if not daytime6thFlr.empty:
                    text_frame = table.rows[currentRow].cells[0].text_frame
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = "6th Floor"
                    font = run.font
                    font.size = Pt(int(daytimeFontSize))
                    font.name = "Arial"
                    font.bold = True
                    font.underline = True
                    font.color.rgb = RGBColor(0xFF, 0xFF, 0x00)

                    text_frame = table.rows[currentRow].cells[4].text_frame
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = "Room"
                    font = run.font
                    font.size = Pt(int(daytimeFontSize))
                    font.name = "Arial"
                    font.bold = True
                    font.underline = True
                    font.color.rgb = RGBColor(0xFF, 0xFF, 0x00)

                    currentRow += 1
                    colorToggle = True

                    for idx, row in daytime6thFlr.iterrows():
                        if colorToggle:
                            color = RGBColor(0xFF, 0xFF, 0xFF)
                        else:
                            color = RGBColor(0xFA, 0xC0, 0x90)
                        colorToggle = not colorToggle
                        for col, value in enumerate(
                            [
                                "Start Time",
                                "End Time",
                                "Section Number",
                                "Section Title",
                                "Room",
                            ]
                        ):
                            text_frame = table.rows[currentRow].cells[col].text_frame
                            p = text_frame.paragraphs[0]
                            run = p.add_run()
                            if value == "Start Time" or value == "End Time":
                                run.text = str(row[value]).lstrip("0")
                            elif value == "Room":
                                run.text = str(row[value]).lstrip("Classroom")
                            else:
                                run.text = str(row[value])
                            font = run.font
                            font.size = Pt(int(daytimeFontSize))
                            font.name = "Arial"
                            font.bold = True
                            font.color.rgb = color
                        currentRow += 1
                    currentRow += 1

                if not daytime7thFlr.empty:
                    text_frame = table.rows[currentRow].cells[0].text_frame
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = "7th Floor"
                    font = run.font
                    font.size = Pt(int(daytimeFontSize))
                    font.name = "Arial"
                    font.bold = True
                    font.underline = True
                    font.color.rgb = RGBColor(0xFF, 0xFF, 0x00)

                    text_frame = table.rows[currentRow].cells[4].text_frame
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = "Room"
                    font = run.font
                    font.size = Pt(int(daytimeFontSize))
                    font.name = "Arial"
                    font.bold = True
                    font.underline = True
                    font.color.rgb = RGBColor(0xFF, 0xFF, 0x00)

                    currentRow += 1
                    colorToggle = True

                    for idx, row in daytime7thFlr.iterrows():
                        if colorToggle:
                            color = RGBColor(0xFF, 0xFF, 0xFF)
                        else:
                            color = RGBColor(0xFA, 0xC0, 0x90)
                        colorToggle = not colorToggle
                        for col, value in enumerate(
                            [
                                "Start Time",
                                "End Time",
                                "Section Number",
                                "Section Title",
                                "Room",
                            ]
                        ):
                            text_frame = table.rows[currentRow].cells[col].text_frame
                            p = text_frame.paragraphs[0]
                            run = p.add_run()
                            if value == "Start Time" or value == "End Time":
                                run.text = str(row[value]).lstrip("0")
                            elif value == "Room":
                                run.text = str(row[value]).lstrip("Classroom")
                            else:
                                run.text = str(row[value])
                            font = run.font
                            font.size = Pt(int(daytimeFontSize))
                            font.name = "Arial"
                            font.bold = True
                            font.color.rgb = color
                        currentRow += 1
            else:
                slide._element.set("show", "0")

            slide = prs.slides[2]
            if not eveningBlock.empty:
                text_frame = slide.shapes[2].text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = (
                    f"UC Berkeley Extension - {dateList[i].strftime('%A')}, "
                    f"{dateList[i].strftime('%B %d, %Y').replace(' 0', ' ')}"
                )
                font = run.font
                font.size = Pt(70)
                font.name = "Calibri"
                font.bold = True
                font.color.rgb = RGBColor(0xFF, 0xFF, 0x00)
                table = slide.shapes[0].table

                currentRow = 0
                if not evening5thFlr.empty:
                    text_frame = table.rows[currentRow].cells[0].text_frame
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = "5th Floor"
                    font = run.font
                    font.size = Pt(int(eveningFontSize))
                    font.name = "Arial"
                    font.bold = True
                    font.underline = True
                    font.color.rgb = RGBColor(0xFF, 0xFF, 0x00)

                    text_frame = table.rows[currentRow].cells[4].text_frame
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = "Room"
                    font = run.font
                    font.size = Pt(int(eveningFontSize))
                    font.name = "Arial"
                    font.bold = True
                    font.underline = True
                    font.color.rgb = RGBColor(0xFF, 0xFF, 0x00)

                    currentRow += 1
                    colorToggle = True

                    for idx, row in evening5thFlr.iterrows():
                        if colorToggle:
                            color = RGBColor(0xFF, 0xFF, 0xFF)
                        else:
                            color = RGBColor(0xFA, 0xC0, 0x90)
                        colorToggle = not colorToggle
                        for col, value in enumerate(
                            [
                                "Start Time",
                                "End Time",
                                "Section Number",
                                "Section Title",
                                "Room",
                            ]
                        ):
                            text_frame = table.rows[currentRow].cells[col].text_frame
                            p = text_frame.paragraphs[0]
                            run = p.add_run()
                            if value == "Start Time" or value == "End Time":
                                run.text = str(row[value]).lstrip("0")
                            elif value == "Room":
                                run.text = str(row[value]).lstrip("Classroom")
                            else:
                                run.text = str(row[value])
                            font = run.font
                            font.size = Pt(int(eveningFontSize))
                            font.name = "Arial"
                            font.bold = True
                            font.color.rgb = color
                        currentRow += 1
                    currentRow += 1

                if not evening6thFlr.empty:
                    text_frame = table.rows[currentRow].cells[0].text_frame
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = "6th Floor"
                    font = run.font
                    font.size = Pt(int(eveningFontSize))
                    font.name = "Arial"
                    font.bold = True
                    font.underline = True
                    font.color.rgb = RGBColor(0xFF, 0xFF, 0x00)

                    text_frame = table.rows[currentRow].cells[4].text_frame
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = "Room"
                    font = run.font
                    font.size = Pt(int(eveningFontSize))
                    font.name = "Arial"
                    font.bold = True
                    font.underline = True
                    font.color.rgb = RGBColor(0xFF, 0xFF, 0x00)

                    currentRow += 1
                    colorToggle = True

                    for idx, row in evening6thFlr.iterrows():
                        if colorToggle:
                            color = RGBColor(0xFF, 0xFF, 0xFF)
                        else:
                            color = RGBColor(0xFA, 0xC0, 0x90)
                        colorToggle = not colorToggle
                        for col, value in enumerate(
                            [
                                "Start Time",
                                "End Time",
                                "Section Number",
                                "Section Title",
                                "Room",
                            ]
                        ):
                            text_frame = table.rows[currentRow].cells[col].text_frame
                            p = text_frame.paragraphs[0]
                            run = p.add_run()
                            if value == "Start Time" or value == "End Time":
                                run.text = str(row[value]).lstrip("0")
                            elif value == "Room":
                                run.text = str(row[value]).lstrip("Classroom")
                            else:
                                run.text = str(row[value])
                            font = run.font
                            font.size = Pt(int(eveningFontSize))
                            font.name = "Arial"
                            font.bold = True
                            font.color.rgb = color
                        currentRow += 1
                    currentRow += 1

                if not evening7thFlr.empty:
                    text_frame = table.rows[currentRow].cells[0].text_frame
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = "7th Floor"
                    font = run.font
                    font.size = Pt(int(eveningFontSize))
                    font.name = "Arial"
                    font.bold = True
                    font.underline = True
                    font.color.rgb = RGBColor(0xFF, 0xFF, 0x00)

                    text_frame = table.rows[currentRow].cells[4].text_frame
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = "Room"
                    font = run.font
                    font.size = Pt(int(eveningFontSize))
                    font.name = "Arial"
                    font.bold = True
                    font.underline = True
                    font.color.rgb = RGBColor(0xFF, 0xFF, 0x00)

                    currentRow += 1
                    colorToggle = True

                    for idx, row in evening7thFlr.iterrows():
                        if colorToggle:
                            color = RGBColor(0xFF, 0xFF, 0xFF)
                        else:
                            color = RGBColor(0xFA, 0xC0, 0x90)
                        colorToggle = not colorToggle
                        for col, value in enumerate(
                            [
                                "Start Time",
                                "End Time",
                                "Section Number",
                                "Section Title",
                                "Room",
                            ]
                        ):
                            text_frame = table.rows[currentRow].cells[col].text_frame
                            p = text_frame.paragraphs[0]
                            run = p.add_run()
                            if value == "Start Time" or value == "End Time":
                                run.text = str(row[value]).lstrip("0")
                            elif value == "Room":
                                run.text = str(row[value]).lstrip("Classroom")
                            else:
                                run.text = str(row[value])
                            font = run.font
                            font.size = Pt(int(eveningFontSize))
                            font.name = "Arial"
                            font.bold = True
                            font.color.rgb = color
                        currentRow += 1
            else:
                slide._element.set("show", "0")

            # Save as Microsoft Powerpoint pptx file, per day
            prs.save(
                (
                    f"{self.saveSignsDirectory}\\{location} "
                    f"{dateList[i].strftime('%Y-%m-%d')} "
                    f"{dateList[i].strftime('%A')}.pptx"
                )
            )
        return 1

    def SFCScheduleToGSheets(
        self, date: datetime.datetime, schedule: pd.DataFrame
    ) -> pygsheets.PyGsheetsException:
        # Sort schedule by time of day blocks and floor
        daytimeBlock = (
            schedule.set_index("Start Time")
            .between_time("00:00", "17:00", inclusive="left")
            .reset_index()
        )
        daytimeBlock = daytimeBlock.sort_values(by=["Room", "Start Time"])
        daytimeBlock["Start Time"] = daytimeBlock["Start Time"].dt.strftime("%I:%M %p")
        daytimeBlock["End Time"] = daytimeBlock["End Time"].dt.strftime("%I:%M %p")
        daytime5thFlr = daytimeBlock.loc[daytimeBlock["Room"] <= "Classroom 515", :]
        daytime6thFlr = daytimeBlock.loc[
            (daytimeBlock["Room"] >= "Classroom 602")
            & (daytimeBlock["Room"] <= "Classroom 613"),
            :,
        ]
        daytime7thFlr = daytimeBlock.loc[daytimeBlock["Room"] >= "Classroom 702", :]

        eveningBlock = (
            schedule.set_index("Start Time")
            .between_time("17:00", "00:00", inclusive="left")
            .reset_index()
        )
        eveningBlock["Start Time"] = eveningBlock["Start Time"].dt.strftime("%I:%M %p")
        eveningBlock["End Time"] = eveningBlock["End Time"].dt.strftime("%I:%M %p")
        eveningBlock = eveningBlock.sort_values(by=["Room", "Start Time"])
        evening5thFlr = eveningBlock.loc[eveningBlock["Room"] <= "Classroom 515", :]
        evening6thFlr = eveningBlock.loc[
            (eveningBlock["Room"] >= "Classroom 602")
            & (eveningBlock["Room"] <= "Classroom 613"),
            :,
        ]
        evening7thFlr = eveningBlock.loc[eveningBlock["Room"] >= "Classroom 702", :]

        blockList = ["Daytime", "Evening"]
        floorList = [
            (
                ("5th Floor", daytime5thFlr),
                ("6th Floor", daytime6thFlr),
                ("7th Floor", daytime7thFlr),
            ),
            (
                ("5th Floor", evening5thFlr),
                ("6th Floor", evening6thFlr),
                ("7th Floor", evening7thFlr),
            ),
        ]

        # Connect to Google Sheets and update with current schedule
        dirpath = os.getcwd()
        client = pygsheets.authorize(service_file=f"{dirpath}/service_file.json")
        try:
            sheet = client.open_by_url(self.SFCScheduleURL)
        except pygsheets.SpreadsheetNotFound as error:
            print(error)
            return error
        finally:
            for i in range(len(blockList)):
                wks = sheet.worksheet_by_title(blockList[i])
                wks.clear(start="A1", end=None, fields="*")
                wks.resize(len(schedule.index) + 6, 5)
                wks.update_value(
                    "A1",
                    (
                        f"UC Berkeley Extension - {date.strftime('%A')}, "
                        f"{date.strftime('%B %d, %Y').replace(' 0', ' ')}"
                    ),
                )
                wks.update_row(
                    2, ["Start Time", "End Time", "Section Title", "Instructor", "Room"]
                )
                rowNumber = 3
                for floor in floorList[i]:
                    if not floor[1].empty:
                        wks.update_value(f"A{rowNumber}", floor[0])
                        rowNumber += 1
                        wks.set_dataframe(
                            floor[1][
                                [
                                    "Start Time",
                                    "End Time",
                                    "Section Title",
                                    "Instructor",
                                    "Room",
                                ]
                            ],
                            f"A{rowNumber}",
                            copy_head=False,
                            fit=False,
                        )
                        rowNumber += len(floor[1].index) + 1


if __name__ == "__main__":
    # os.environ["QT_AUTO_SCREEN_FACTOR"] = "1"
    app = QtWidgets.QApplication(sys.argv)
    mainWindow = QtWidgets.QWidget()
    ui = Ui_mainWindow()
    ui.setupUi(mainWindow)
    mainWindow.show()
    sys.exit(app.exec_())
